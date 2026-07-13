# =============================================================================
# SYRA AI PLATFORM — Production Backend
# =============================================================================

import os
import re
import json
import shutil
import secrets
import traceback
import zipfile
import uuid
import base64
import hmac
from datetime import datetime, timedelta

import pandas as pd
import requests
from bs4 import BeautifulSoup
from flask import (
    Flask, request, jsonify, session,
    send_from_directory, Response, stream_with_context
)
from flask_cors import CORS
from werkzeug.utils import secure_filename
from werkzeug.security import generate_password_hash, check_password_hash

import replicate
import random
import smtplib
import ssl
import sqlite3
import hashlib
from email.message import EmailMessage
from datetime import datetime, timedelta
from dotenv import load_dotenv

import firebase_admin
from firebase_admin import credentials, auth as firebase_auth

# Centralized AI routing with multi-provider fallback and category-aware routing
# (General Chat, Coding, Website Generation, Project Generation, Agriculture,
# and streaming chat all go through this module — see ai_engine.py for the
# single source of truth on providers/models and fallback chains).
from ai_engine import generate_ai_reply, generate_ai_reply_stream, FRIENDLY_ERROR, get_provider_chain, cloudflare_speech_to_text

load_dotenv()

print("Replicate Token:", os.getenv("REPLICATE_API_TOKEN"))

# =============================================================================
# ENVIRONMENT & API KEYS
# =============================================================================
OPENROUTER_KEY = os.environ.get("OPENROUTER_API_KEY", "")
OPENROUTER_URL = "https://openrouter.ai/api/v1/chat/completions"

REPLICATE_TOKEN = os.environ.get("REPLICATE_API_TOKEN", "")
if REPLICATE_TOKEN:
    os.environ["REPLICATE_API_TOKEN"] = REPLICATE_TOKEN

# Cloudflare Workers AI — new PRIMARY image generation provider (Replicate
# above remains as the fallback if Cloudflare is unavailable or its daily
# free quota is exhausted). Purely additive: nothing here touches the
# existing Replicate config/usage.
CLOUDFLARE_ACCOUNT_ID = os.environ.get("CLOUDFLARE_ACCOUNT_ID", "")
CLOUDFLARE_API_TOKEN = os.environ.get("CLOUDFLARE_API_TOKEN", "")
CLOUDFLARE_IMAGE_MODEL = os.environ.get(
    "CLOUDFLARE_IMAGE_MODEL", "@cf/black-forest-labs/flux-1-schnell"
)
if not (CLOUDFLARE_ACCOUNT_ID and CLOUDFLARE_API_TOKEN):
    print("WARNING: No CLOUDFLARE_ACCOUNT_ID/CLOUDFLARE_API_TOKEN set in environment. "
          "Image generation will skip Cloudflare and go straight to the Replicate fallback.")

# SECURITY: never bake a real key in as a source-code default — anyone with
# read access to this file (or a public repo/history) gets your key. Set
# GEMINI_API_KEY in the environment; there is intentionally no fallback value.
GEMINI_API_KEY = os.environ.get("GEMINI_API_KEY", "")

# Fail loudly if critical keys are missing, instead of silently running
# with no AI provider available.
if not OPENROUTER_KEY and not GEMINI_API_KEY:
    print("WARNING: No OPENROUTER_API_KEY or GEMINI_API_KEY set in environment. "
          "Chat features will not work until at least one is configured.")
if not REPLICATE_TOKEN:
    print("WARNING: No REPLICATE_API_TOKEN set in environment. "
          "Image generation and voice STT will not work until it is configured.")

# =============================================================================
# FLASK APP INIT
# =============================================================================
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
PUBLIC_FOLDER = os.path.join(BASE_DIR, "..", "Frontend", "public")
UPLOAD_FOLDER = "./uploaded_datasets"
PROJECTS_FOLDER = "./generated_projects"
MEMORY_FOLDER = "./memory"
VOICES_FOLDER = "./voices"
AUTH_FOLDER = "./auth_data"
AVATAR_FOLDER = "./avatars"

for folder in [PUBLIC_FOLDER, UPLOAD_FOLDER, PROJECTS_FOLDER, MEMORY_FOLDER, VOICES_FOLDER, AUTH_FOLDER, AVATAR_FOLDER]:
    os.makedirs(folder, exist_ok=True)

# =============================================================================
# FIREBASE ADMIN SDK INIT (Google Sign-In) — initialized exactly once
# =============================================================================
import json

if not firebase_admin._apps:
    try:
        firebase_json = os.environ.get("FIREBASE_SERVICE_ACCOUNT")

        if firebase_json:
            cred_dict = json.loads(firebase_json)
            firebase_admin.initialize_app(
                credentials.Certificate(cred_dict)
            )
            print("Firebase Admin SDK initialized from environment.")
        else:
            print("WARNING: FIREBASE_SERVICE_ACCOUNT environment variable not found.")

    except Exception as e:
        print(f"WARNING: Firebase initialization failed: {e}")

app = Flask(__name__, static_folder=PUBLIC_FOLDER, static_url_path="")
app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config["MAX_CONTENT_LENGTH"] = 100 * 1024 * 1024

app.secret_key = os.environ.get("SECRET_KEY") or os.urandom(32).hex()

# OTP Settings
otp_store = {}

OTP_EXPIRE_MINUTES = 5

DATABASE = "syra.db"

if not os.environ.get("SECRET_KEY"):
    print("WARNING: No SECRET_KEY set in environment. Using a randomly generated "
          "key for this process — user sessions will be invalidated on every "
          "restart. Set SECRET_KEY in your environment for production.")
    
def get_db():
    """Returns a sqlite3 connection with Row access (dict-like columns),
    used by the chat history / profile / settings persistence layer."""
    conn = sqlite3.connect(DATABASE)
    conn.row_factory = sqlite3.Row
    return conn


def init_database():

    conn = sqlite3.connect(DATABASE)
    cur = conn.cursor()

    cur.execute("""
    CREATE TABLE IF NOT EXISTS users(

        id INTEGER PRIMARY KEY AUTOINCREMENT,

        name TEXT,

        email TEXT UNIQUE,

        phone TEXT UNIQUE,

        login_type TEXT,

        google_id TEXT,

        created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP

    )
    """)

    # Backward-compatible column add: existing DBs created before the
    # profile picture feature won't have this column yet.
    try:
        cur.execute("ALTER TABLE users ADD COLUMN picture TEXT")
    except sqlite3.OperationalError:
        pass  # column already exists

    # ---------------------------------------------------------------
    # CHAT HISTORY (SQLite-backed) — replaces the earlier per-user JSON
    # file store. user_key is the same bucketing key resolve_session_id()
    # already produces (logged-in email, or a guest session_id/"anon").
    # ---------------------------------------------------------------
    cur.execute("""
    CREATE TABLE IF NOT EXISTS conversations(
        id TEXT NOT NULL,
        user_key TEXT NOT NULL,
        title TEXT NOT NULL DEFAULT 'New chat',
        pinned INTEGER NOT NULL DEFAULT 0,
        created_at TEXT NOT NULL,
        updated_at TEXT NOT NULL,
        PRIMARY KEY (user_key, id)
    )
    """)

    cur.execute("""
    CREATE TABLE IF NOT EXISTS messages(
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_key TEXT NOT NULL,
        conv_id TEXT NOT NULL,
        role TEXT NOT NULL,
        content TEXT NOT NULL,
        created_at TEXT NOT NULL
    )
    """)
    cur.execute("""
    CREATE INDEX IF NOT EXISTS idx_messages_user_conv
    ON messages(user_key, conv_id)
    """)

    # ---------------------------------------------------------------
    # PER-USER SETTINGS (dark mode, language, theme)
    # ---------------------------------------------------------------
    cur.execute("""
    CREATE TABLE IF NOT EXISTS user_settings(
        user_key TEXT PRIMARY KEY,
        dark_mode INTEGER NOT NULL DEFAULT 1,
        language TEXT NOT NULL DEFAULT 'en',
        theme TEXT NOT NULL DEFAULT 'dark',
        updated_at TEXT
    )
    """)

    conn.commit()
    conn.close()


def send_email_otp(email, otp):

    msg = EmailMessage()

    msg["Subject"] = "SYRA Login OTP"

    msg["From"] = os.getenv("EMAIL_ADDRESS")

    msg["To"] = email

    msg.set_content(f"""

Your SYRA Verification Code

OTP : {otp}

This OTP expires in 5 minutes.

Do not share this code.

""")

    context = ssl.create_default_context()

    with smtplib.SMTP_SSL(
        "smtp.gmail.com",
        465,
        context=context
    ) as smtp:

        smtp.login(

            os.getenv("EMAIL_ADDRESS"),

            os.getenv("EMAIL_PASSWORD")

        )

        smtp.send_message(msg)

# SECURE SESSION COOKIE CONFIG
# - HttpOnly: JS on the page can never read the session cookie (mitigates XSS
#   cookie theft).
# - SameSite=Lax: cookie isn't sent on most cross-site requests (mitigates CSRF
#   on GET-triggered navigation); combined with the explicit CSRF token check
#   below for state-changing requests.
# - Secure: only sent over HTTPS. Opt-in via COOKIE_SECURE=1 once you're behind
#   HTTPS in production — left off by default so local http://localhost dev
#   still works (browsers drop Secure cookies on plain http).
# - PERMANENT_SESSION_LIFETIME: how long a "Remember Me" session lasts. A
#   non-permanent session (remember unchecked) instead expires when the
#   browser closes.
app.config["SESSION_COOKIE_HTTPONLY"] = True
app.config["SESSION_COOKIE_SAMESITE"] = "None"
app.config["SESSION_COOKIE_SECURE"] = True
app.config["PERMANENT_SESSION_LIFETIME"] = timedelta(days=30)

# CORS: a wildcard origin combined with credentialed requests (cookies) is
# unsafe and rejected by most browsers anyway. Configure real origins via
# ALLOWED_ORIGINS env var (comma-separated). Defaults to localhost for dev.
_allowed_origins = os.environ.get("ALLOWED_ORIGINS", "http://localhost:5000,http://127.0.0.1:5000,https://syra-ai.netlify.app")
CORS(app, supports_credentials=True, origins=[o.strip() for o in _allowed_origins.split(",") if o.strip()])

# =============================================================================
# CONSTANTS & CONFIGURATION
# =============================================================================
ALLOWED_EXTENSIONS = {
    # documents / data
    "csv", "txt", "pdf", "docx", "doc", "xlsx", "xls", "pptx", "json", "xml", "md",
    # images
    "png", "jpg", "jpeg", "gif", "webp",
    # audio
    "mp3", "wav", "ogg",
    # archives
    "zip",
    # code
    "py", "html", "htm", "css", "js", "jsx", "ts", "tsx", "java", "c", "cpp", "h",
    "cs", "php", "go", "rs",
}
IMAGE_EXTENSIONS = {"png", "jpg", "jpeg", "gif", "webp"}
CODE_EXTENSIONS = {"py", "html", "htm", "css", "js", "jsx", "ts", "tsx", "java", "c",
                    "cpp", "h", "cs", "php", "go", "rs"}
TOP_K = 6
SCORE_THRESHOLD = 0.75
TEXT_COL = "_combined_text"
MAX_HISTORY = 30

# NOTE: model names/providers used to be hardcoded here (MODEL_MAP) and
# duplicated inside call_ai()/call_ai_stream(). That routing now lives
# entirely in ai_engine.py — see generate_ai_reply()/generate_ai_reply_stream()
# below, which is the single place every AI feature (general chat, coding,
# website generation, project generation) goes through.

# =============================================================================
# LANGUAGE & INTENT DETECTION
# =============================================================================
def detect_language(text):
    for ch in text:
        if "\u0B80" <= ch <= "\u0BFF":
            return "ta"
    tanglish = ["enna", "epdi", "iruku", "venum", "yaar", "machan", "machi",
                "poda", "vada", "seri", "sari", "nala", "romba", "evalo", "evlo"]
    if any(w in text.lower() for w in tanglish):
        return "ta_tanglish"
    return "en"

def detect_intent(text):
    t = text.lower()
    if any(k in t for k in ["crop", "fertilizer", "irrigation", "farming", "agriculture", "payir", "nel"]):
        return "agri"
    if any(k in t for k in ["image", "photo", "picture", "draw", "generate image", "create image"]):
        return "image"
    if any(k in t for k in ["build website", "create website", "make website", "generate website", "html", "css", "react"]):
        return "website"
    if any(k in t for k in ["build project", "create project", "generate project", "full project", "ecommerce", "portfolio"]):
        return "project"
    if any(k in t for k in ["write code", "debug", "fix code", "explain code", "refactor"]):
        return "code"
    return "chat"

# =============================================================================
# SMART LIVE SEARCH DETECTION — triggers automatic web grounding
# =============================================================================
# Comprehensive keyword map organized by topic. ALL groups are checked; a
# match in ANY group triggers live search. Purely additive — only ever turns
# search ON, never off, and never replaces the manual toggle.
_LIVE_DATA_KEYWORDS = {
    # ------- Sports -------
    "sports": [
        "score", "scores", "match", "fixture", "fixtures", "standings", "live score",
        "football", "cricket", "ipl", "nba", "formula 1", "f1", "olympics",
        "kabaddi", "tennis", "tournament", "world cup", "premier league",
        "test match", "t20", "odi", "wicket", "goal", "league table",
        "champions league", "series result", "squad", "playing xi",
        "basketball", "rugby", "hockey", "baseball", "golf tournament",
        "race result", "grand prix", "sprint race",
    ],
    # ------- Breaking news -------
    "news": [
        "latest news", "breaking news", "news today", "headlines", "current affairs",
        "today's news", "what happened", "recent news", "news update",
        "just happened", "happened today", "this week news", "this month news",
        "latest update", "what's happening",
    ],
    # ------- Weather -------
    "weather": [
        "weather", "forecast", "temperature today", "rain today", "humidity",
        "wind speed", "uv index", "feels like", "tomorrow weather",
        "weekly forecast", "monsoon update", "cyclone", "flood alert",
        "heatwave", "snowfall", "drizzle today",
    ],
    # ------- Finance & markets -------
    "market": [
        "stock price", "share price", "stock market", "nifty", "sensex",
        "gold price", "crypto price", "bitcoin price", "exchange rate",
        "currency rate", "market today", "dow jones", "nasdaq", "s&p 500",
        "ethereum price", "silver price", "crude oil price", "forex",
        "interest rate", "rbi rate", "fed rate", "inflation rate",
        "market cap", "ipo today", "share market", "petrol price",
        "diesel price", "fuel price", "commodity price",
    ],
    # ------- Politics & government leaders -------
    "politics": [
        "prime minister", "chief minister", "president", "governor",
        "current pm", "current cm", "who is the pm", "who is the cm",
        "who is president", "minister of", "cabinet minister",
        "election result", "who won election", "elected leader",
        "mp of", "mla of", "lok sabha", "rajya sabha", "parliament",
        "government policy", "new law", "policy announced",
        "budget 2024", "budget 2025", "budget 2026",
        "opposition leader", "ruling party", "political news",
        "chief justice", "speaker of", "who is the mayor",
        "senator", "congress member", "mp result",
    ],
    # ------- Technology & AI -------
    "technology": [
        "latest ai", "new ai model", "chatgpt update", "gemini update",
        "claude update", "gpt-5", "gpt 5", "llm news", "ai news",
        "new phone", "new iphone", "new android", "new laptop",
        "software update", "new version", "release notes", "launched today",
        "tech news", "startup news", "product launch", "new release",
        "latest version of", "current version of", "update available",
        "apple event", "google io", "microsoft event",
    ],
    # ------- Science & health -------
    "health": [
        "outbreak", "epidemic", "pandemic", "vaccination", "vaccine update",
        "new drug", "clinical trial", "health advisory", "who advisory",
        "disease spread", "treatment update", "health news",
        "new study", "research finding", "new discovery",
    ],
    # ------- Direct live-data language -------
    "live_queries": [
        "right now", "at this moment", "currently", "live",
        "today's", "this year's", "what is today",
        "latest", "most recent", "just released", "just announced",
        "who won", "who is winning", "who leads",
        "how much is", "what is the price of", "what is the rate of",
        "what is the current", "who is currently", "who is the current",
    ],
}

# Keywords that indicate a forward-looking prediction/forecast query.
# When detected, the AI is instructed to label its answer as a prediction,
# not a verified fact, before proceeding.
_PREDICTION_KEYWORDS = [
    "predict", "prediction", "forecast", "estimate", "projection",
    "will happen", "what will", "who will win", "who will be",
    "will it", "future of", "trend", "next year", "next month",
    "upcoming election", "likely to", "expected to", "chances of",
    "probability", "odds of", "by 2025", "by 2026", "by 2030",
    "prognosis", "outlook", "going to be", "when will",
]


def needs_live_data(text):
    """Returns True if the query requires live web grounding.
    Triggers automatically for: sports results, breaking news, weather,
    financial markets, political leaders, technology launches, health
    advisories, and any query using live-data language like 'currently',
    'latest', or 'who is the current ...'."""
    t = text.lower()
    for keywords in _LIVE_DATA_KEYWORDS.values():
        if any(k in t for k in keywords):
            return True
    return False


def detect_prediction_mode(text):
    """Returns True when the user is asking for a future prediction,
    forecast, or estimate (not a verified current fact). The model will
    explicitly label its response as a prediction rather than fact."""
    t = text.lower()
    return any(k in t for k in _PREDICTION_KEYWORDS)

def map_intent_to_category(intent):
    """Maps detected intent to AI engine category for smart provider routing."""
    intent_map = {
        "agri": "agri",
        "image": "image",  # Not routed to AI engine
        "website": "website",
        "project": "project",
        "code": "coding",
        "chat": "chat",
    }
    return intent_map.get(intent, "chat")

def handle_identity_questions(text):
    t = text.lower()
    if "your name" in t or "who are you" in t:
        return "My name is SYRA. I am an advanced, context-aware AI assistant created by Surya under KRISH."
    if "founder" in t or "who created you" in t:
        return "I was created by Surya, the founder of KRISH."
    if "company" in t:
        return "I was developed under the company KRISH."
    return None

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

def current_date_context():
    """Builds a fresh 'today is...' system snippet from the server clock so
    SYRA always knows the real current date/time without needing web search,
    instead of relying on whatever date is baked into training data."""
    now = datetime.now()
    return (
        f"\n\nCurrent date and time (server clock, authoritative — always use "
        f"this instead of any date you might otherwise assume): "
        f"{now.strftime('%A, %B %d, %Y, %I:%M %p')}."
    )

# =============================================================================
# LIVE WEB SEARCH — multi-source verification, trust-ranked, conflict- and
# confidence-aware grounding pipeline (Perplexity/ChatGPT-Search-style)
# =============================================================================
import difflib
import concurrent.futures
from collections import defaultdict

# ---------------------------------------------------------------------------
# SOURCE PRIORITY
# ---------------------------------------------------------------------------
# Domains grouped into trust tiers, lowest index = highest trust:
#   0 Government            1 Official orgs / institutions / sports bodies
#   2 Official documentation 3 Wikipedia    4 Official company sites
#   5 Academic / research    6 Trusted news / finance data
#   7 Community / discussion (still usable, lowest trust tier)
# Anything not listed is still used (never silently dropped) but ranks last.
_DOMAIN_TIERS = [
    # Tier 0 — Government / official regulators
    {
        "gov.in", "gov.uk", "usa.gov", "nic.in", "eci.gov.in", "parliament.uk",
        "parliament.in", "irs.gov", "europa.eu", "cdc.gov", "mohfw.gov.in",
        "nih.gov", "rbi.org.in", "sebi.gov.in", "whitehouse.gov",
    },
    # Tier 1 — Official organizations / international institutions / official
    # sports governing bodies
    {
        "who.int", "un.org", "imf.org", "worldbank.org", "unicef.org",
        "fifa.com", "icc-cricket.com", "bcci.tv", "nba.com", "formula1.com",
        "olympics.com", "uefa.com",
    },
    # Tier 2 — Official documentation
    {
        "python.org", "docs.python.org", "developer.mozilla.org",
        "reactjs.org", "react.dev", "nodejs.org", "learn.microsoft.com",
        "docs.microsoft.com", "cloud.google.com", "developer.apple.com",
        "developers.google.com",
    },
    # Tier 3 — Wikipedia
    {"wikipedia.org"},
    # Tier 4 — Official company websites
    {
        "openai.com", "anthropic.com", "blog.google", "apple.com",
        "microsoft.com", "google.com", "meta.com",
    },
    # Tier 5 — Academic / research sources
    {
        "nature.com", "sciencedirect.com", "springer.com", "arxiv.org",
        "ncbi.nlm.nih.gov", "jstor.org",
    },
    # Tier 6 — Trusted news & financial data outlets
    {
        "reuters.com", "bbc.com", "bbc.co.uk", "theguardian.com", "apnews.com",
        "thehindu.com", "ndtv.com", "hindustantimes.com", "timesofindia.com",
        "economictimes.com", "livemint.com", "businessstandard.com",
        "bloomberg.com", "ft.com", "wsj.com", "marketwatch.com",
        "moneycontrol.com", "investing.com", "nseindia.com", "bseindia.com",
        "espn.com", "espncricinfo.com", "cricbuzz.com", "goal.com",
        "techcrunch.com", "theverge.com", "arstechnica.com", "wired.com",
        "britannica.com",
    },
    # Tier 7 — Community / discussion
    {"stackoverflow.com", "reddit.com", "quora.com", "medium.com"},
]

# Flat set of every recognised trusted domain. Kept under the original name
# so any other code (e.g. the /api/search endpoint) that references
# _TRUSTED_DOMAINS keeps working unchanged.
_TRUSTED_DOMAINS = frozenset(d for tier in _DOMAIN_TIERS for d in tier)

# Category-specific domains promoted ABOVE every general tier (rank -1) when
# the query matches that category — implements "programming questions
# prioritise Python Docs/MDN/...", "sports prioritises ICC/BCCI/...", and
# "finance prioritises official financial data".
_CATEGORY_PRIORITY_DOMAINS = {
    "programming": {
        "docs.python.org", "python.org", "developer.mozilla.org",
        "reactjs.org", "react.dev", "nodejs.org", "learn.microsoft.com",
        "docs.microsoft.com", "cloud.google.com", "developers.google.com",
        "stackoverflow.com",
    },
    "sports": {
        "icc-cricket.com", "bcci.tv", "espncricinfo.com", "cricbuzz.com",
        "fifa.com", "nba.com", "formula1.com", "espn.com", "olympics.com",
        "uefa.com",
    },
    "finance": {
        "rbi.org.in", "sebi.gov.in", "nseindia.com", "bseindia.com",
        "moneycontrol.com", "bloomberg.com", "wsj.com", "marketwatch.com",
        "investing.com", "ft.com",
    },
}

_PROGRAMMING_KEYWORDS = [
    "python", "javascript", "typescript", "react", "node.js", "nodejs",
    "flask", "django", "api reference", "documentation", "syntax error",
    "stack trace", "npm install", "pip install", "library version",
    "package version", "sdk", "framework docs", "error code", "compiler",
]


def _detect_search_category(text):
    """Best-effort topic category used purely for source-priority weighting.
    Reuses the same signal words as needs_live_data() so behaviour stays
    consistent. Returns 'programming' | 'sports' | 'finance' | 'general'."""
    t = text.lower()
    if any(k in t for k in _PROGRAMMING_KEYWORDS):
        return "programming"
    if any(k in t for k in _LIVE_DATA_KEYWORDS["sports"]):
        return "sports"
    if any(k in t for k in _LIVE_DATA_KEYWORDS["market"]):
        return "finance"
    return "general"


def _domain_rank(url: str, category: str = "general") -> int:
    """Lower integer = higher trust/priority. Category-priority domains rank
    highest (-1), then general trust tiers 0..7, then unranked domains last
    (still included — never silently dropped)."""
    priority_set = _CATEGORY_PRIORITY_DOMAINS.get(category)
    if priority_set and any(d in url for d in priority_set):
        return -1
    for tier_index, tier_domains in enumerate(_DOMAIN_TIERS):
        if any(d in url for d in tier_domains):
            return tier_index
    return len(_DOMAIN_TIERS)


def _clean_ddg_url(href: str) -> str:
    """DuckDuckGo wraps outbound links as /l/?uddg=<encoded-url>.
    Decode to the real destination URL so callers can fetch the page."""
    if not href:
        return ""
    if href.startswith("/l/?") or "duckduckgo.com/l/?" in href:
        try:
            from urllib.parse import urlparse, parse_qs, unquote
            qs = parse_qs(urlparse(href).query)
            uddg = qs.get("uddg", [""])[0]
            if uddg:
                return unquote(uddg)
        except Exception:
            pass
    return href


# Simple in-process TTL caches.
# _search_cache: query+limit+category -> (timestamp, results_list) — 2 min.
# _page_cache: url -> (timestamp, extracted_text) — 5 min, shared across
# queries so re-grounding the same trending source doesn't re-scrape it.
_search_cache: dict = {}
_CACHE_TTL_SECONDS = 120
_page_cache: dict = {}
_PAGE_CACHE_TTL_SECONDS = 300


def web_search(query, limit=5, category="general"):
    """Returns a list of {title, snippet, url} dicts ordered by source
    priority (category-aware), deduplicated by URL. Returns [] on failure —
    never injects fake error strings as search results (those would be
    repeated back as fact by the model)."""
    from urllib.parse import quote as urlquote

    cache_key = f"{query.lower().strip()}:::{limit}:::{category}"
    now = datetime.now().timestamp()
    if cache_key in _search_cache:
        ts, cached = _search_cache[cache_key]
        if now - ts < _CACHE_TTL_SECONDS:
            return cached

    last_err = None
    for attempt in range(2):
        try:
            url = f"https://html.duckduckgo.com/html/?q={urlquote(query)}"
            res = requests.get(
                url, timeout=7,
                headers={"User-Agent": "Mozilla/5.0 (compatible; SYRA-AI/2.0)"}
            )
            if res.status_code != 200:
                last_err = f"status {res.status_code}"
                continue

            soup = BeautifulSoup(res.text, "html.parser")
            seen_urls: set = set()
            raw: list = []

            for r in soup.find_all("div", class_="result"):
                title_el = r.find("a", class_="result__a")
                snippet_el = r.find("a", class_="result__snippet")
                title = title_el.text.strip() if title_el else ""
                snippet = snippet_el.text.strip() if snippet_el else ""
                href = _clean_ddg_url(
                    title_el["href"] if title_el and title_el.get("href") else ""
                )

                if not title or not href:
                    continue

                norm = href.rstrip("/").lower()
                if norm in seen_urls:
                    continue
                seen_urls.add(norm)

                raw.append({
                    "title": title,
                    "snippet": snippet,
                    "url": href,
                    "_rank": _domain_rank(href, category),
                })

                if len(raw) >= limit * 4:   # collect extras for verification, then trim
                    break

            # Stable sort: highest-priority sources first, DDG relevance
            # preserved within each tier — this is the multi-source pool
            # later cross-checked for duplicates/conflicts.
            raw.sort(key=lambda x: x["_rank"])
            out = [{"title": r["title"], "snippet": r["snippet"], "url": r["url"]}
                   for r in raw[:limit]]

            _search_cache[cache_key] = (now, out)
            return out

        except Exception as e:
            last_err = str(e)

    print("web_search failed after retries:", last_err)
    _search_cache[cache_key] = (now, [])
    return []


def summarize_url(url):
    """Fetches a URL and extracts clean, meaningful body text.
    Strips nav / footer / sidebar / scripts so the model gets actual
    article content rather than boilerplate. Cached per-URL for
    _PAGE_CACHE_TTL_SECONDS so repeated grounding of the same trending
    source doesn't re-scrape it every request."""
    if not url or not url.startswith("http"):
        return "Could not fetch URL: invalid URL"

    cache_key = url.strip()
    now = datetime.now().timestamp()
    cached = _page_cache.get(cache_key)
    if cached and now - cached[0] < _PAGE_CACHE_TTL_SECONDS:
        return cached[1]

    try:
        res = requests.get(
            url, timeout=8,
            headers={"User-Agent": "Mozilla/5.0 (compatible; SYRA-AI/2.0)"},
            allow_redirects=True,
        )
        if res.status_code != 200:
            return f"Could not fetch URL: HTTP {res.status_code}"

        soup = BeautifulSoup(res.text, "html.parser")

        # Remove boilerplate
        for tag in soup(["script", "style", "nav", "footer", "header",
                          "aside", "form", "noscript", "iframe",
                          "figure", "figcaption"]):
            tag.decompose()

        # Prefer semantic article / main content sections
        main = (
            soup.find("article")
            or soup.find("main")
            or soup.find(id="content")
            or soup.find(class_="content")
            or soup.find(class_="article-body")
            or soup.find(class_="post-content")
        )
        target = main if main else soup

        text = " ".join(target.get_text(separator=" ").split())
        text = text[:4000]
        _page_cache[cache_key] = (now, text)
        return text
    except requests.Timeout:
        return "Could not fetch URL: timeout"
    except Exception as e:
        return f"Could not fetch URL: {str(e)}"


def wikipedia_lookup(query):
    """Additional trusted knowledge source used alongside DuckDuckGo web_search.
    Returns a short factual summary (Wikipedia's own condensed extract), or
    None on failure. Only ever summarises — never returns raw article text.
    Used as a Wikipedia fallback when live search produces no results, and
    as supplementary depth context when live search does return results."""
    from urllib.parse import quote as urlquote
    # Strip common question words to improve Wikipedia API matching
    clean = query.strip()
    for prefix in ("who is ", "what is ", "who was ", "what was ", "where is "):
        if clean.lower().startswith(prefix):
            clean = clean[len(prefix):]
            break

    for attempt_query in [clean, query]:
        try:
            res = requests.get(
                "https://en.wikipedia.org/api/rest_v1/page/summary/" + urlquote(attempt_query),
                timeout=6,
                headers={"User-Agent": "SYRA-AI/2.0 (https://syra.ai)"},
            )
            if res.status_code != 200:
                continue
            data = res.json()
            extract = (data.get("extract") or "").strip()
            if not extract:
                continue
            title = data.get("title", attempt_query)
            # Limit to 900 chars for context efficiency
            return f"{title}: {extract[:900]}"
        except Exception as e:
            print("wikipedia_lookup failed:", e)
            continue
    return None


# ---------------------------------------------------------------------------
# MULTI-SOURCE VERIFICATION HELPERS
# ---------------------------------------------------------------------------

def _dedupe_near_identical(results):
    """Collapses results whose snippets are near-identical in wording (e.g.
    a wire story syndicated on multiple domains) even when the URLs differ,
    on top of the exact-URL dedup already applied in web_search(). Keeps the
    higher-priority (earlier) copy."""
    kept = []
    seen_texts = []
    for r in results:
        snip = (r.get("snippet") or "").strip()
        if not snip:
            kept.append(r)
            continue
        is_dupe = any(
            difflib.SequenceMatcher(None, snip, s).ratio() > 0.92
            for s in seen_texts
        )
        if is_dupe:
            continue
        seen_texts.append(snip)
        kept.append(r)
    return kept


# Numeric "claims" worth cross-checking across sources: currency amounts,
# percentages, and score-like "a-b" / "a/b" pairs.
_NUM_CLAIM_PATTERN = re.compile(
    r'(?:[$₹€£]\s?\d[\d,]*\.?\d*)'
    r'|(?:\d[\d,]*\.?\d*\s?(?:%|percent))'
    r'|(?:\b\d{1,3}[-/]\d{1,3}\b)',
    re.IGNORECASE,
)


def _extract_numeric_claims(text):
    """Extracts a set of raw numeric-claim substrings (prices, percentages,
    score-like pairs) from a block of text, for cross-source comparison."""
    if not text:
        return set()
    return {m.strip().lower() for m in _NUM_CLAIM_PATTERN.findall(text)}


def _normalize_num(token):
    """Strips currency/percent symbols and thousands separators, returns a
    float value or None if the token doesn't reduce to a plain number."""
    cleaned = re.sub(r'[^\d.]', '', token)
    if not cleaned:
        return None
    try:
        return float(cleaned)
    except ValueError:
        return None


def _detect_conflicts(per_source_claims):
    """Compares numeric claims (prices, percentages, scores) extracted from
    each independent source. Flags a conflict when at least two different
    sources report meaningfully different values for the same kind of
    figure, with no majority agreement — this is the trigger for 'never
    merge conflicting information, never guess'.

    per_source_claims: list of sets, one set of raw claim strings per source.
    Returns (conflict_detected: bool, detail: str | None).
    """
    buckets = defaultdict(list)  # unit -> [(value_or_raw, source_index), ...]

    for idx, claims in enumerate(per_source_claims):
        for c in claims:
            if re.search(r'[-/]', c) and not re.search(r'[%$₹€£]', c):
                buckets["score"].append((c.strip(), idx))
                continue
            val = _normalize_num(c)
            if val is None:
                continue
            if "%" in c or "percent" in c:
                buckets["percent"].append((val, idx))
            elif re.search(r'[$₹€£]', c):
                buckets["currency"].append((val, idx))

    for unit, items in buckets.items():
        distinct_sources = {idx for _, idx in items}
        if len(distinct_sources) < 2:
            continue  # only one source made this kind of claim — nothing to cross-check
        values = [v for v, _ in items]

        if unit == "score":
            distinct_values = sorted(set(values))
            if len(distinct_values) > 1:
                return True, (
                    f"sources report different scores/figures: {', '.join(distinct_values)}"
                )
        else:
            lo, hi = min(values), max(values)
            if lo == 0:
                continue
            # More than ~2% relative spread between the lowest and highest
            # reported values for the same kind of figure = disagreement.
            if (hi - lo) / lo > 0.02:
                return True, (
                    f"sources disagree on a {unit} figure (values range from "
                    f"{lo:g} to {hi:g})"
                )

    return False, None


def _compute_confidence(results, category, conflict_detected):
    """Internal-only confidence tier: 'high' | 'medium' | 'low'. Never
    exposed as a raw score to the user — only used to pick which guidance
    tag gets embedded in the model's grounding context."""
    if not results:
        return "low"
    if conflict_detected:
        return "low"
    trusted_count = sum(1 for r in results if _domain_rank(r.get("url", ""), category) <= 2)
    total = len(results)
    if trusted_count >= 2 and total >= 3:
        return "high"
    if trusted_count >= 1 or total >= 2:
        return "medium"
    return "low"


def build_search_context(query, web_search_enabled=False, auto_live=False,
                         limit=6, is_prediction=False):
    """Builds a grounded system-prompt context block from multi-source live
    search + Wikipedia, with verification, dedup, conflict detection, and
    confidence tagging baked in.

    Returns a string to be appended to the AI's system message, or "" if
    no grounding is needed. Never injects fabricated error strings as facts.

    Behaviour:
    - If neither web_search_enabled nor auto_live → returns "".
    - Detects a topic category (programming / sports / finance / general)
      to apply the right source-priority ordering.
    - Runs web_search() (cached, deduplicated, category-aware trust-ranked).
    - Collapses near-duplicate snippets across different domains.
    - Enriches the top sources with real page content, fetched concurrently.
    - Cross-checks extracted numeric claims (prices/percentages/scores)
      across independent sources and flags disagreement instead of
      merging or guessing.
    - Computes an internal confidence tier and embeds the matching
      instruction (never the raw score) for the model to follow.
    - Falls back to Wikipedia when web search returns nothing.
    - Wikipedia also supplements live results for factual depth.
    - Appends a prediction disclaimer when is_prediction=True.
    """
    if not (web_search_enabled or auto_live):
        return ""

    category = _detect_search_category(query)
    results = web_search(query, limit=limit, category=category)

    if results:
        results = _dedupe_near_identical(results)

        # Enrich the top 3 sources with real page content, fetched
        # concurrently so multi-source verification doesn't multiply
        # response latency.
        enrich_targets = [(i, r) for i, r in enumerate(results[:3]) if r.get("url")]
        enriched = []  # list of (original_index, result_dict, page_text)
        if enrich_targets:
            with concurrent.futures.ThreadPoolExecutor(max_workers=len(enrich_targets)) as ex:
                future_map = {
                    ex.submit(summarize_url, r["url"]): (i, r)
                    for i, r in enrich_targets
                }
                try:
                    for fut in concurrent.futures.as_completed(future_map, timeout=12):
                        i, r = future_map[fut]
                        try:
                            page_text = fut.result()
                        except Exception:
                            page_text = None
                        if page_text and not str(page_text).startswith("Could not fetch"):
                            enriched.append((i, r, page_text))
                except concurrent.futures.TimeoutError:
                    pass  # use whatever finished in time; never block the reply on a slow source
            enriched.sort(key=lambda tup: tup[0])

        # ---- Multi-source verification: cross-check numeric claims ----
        per_source_claims = [
            _extract_numeric_claims(r.get("snippet", "")) for r in results
        ]
        per_source_claims += [
            _extract_numeric_claims(page_text) for _, _, page_text in enriched
        ]
        conflict_detected, conflict_detail = _detect_conflicts(per_source_claims)
        confidence = _compute_confidence(results, category, conflict_detected)

        snippet_lines = "\n".join(
            f"• {r.get('title', '')} ({r.get('url', '')}): {r.get('snippet', '')}"
            for r in results
            if r.get("title")
        )

        context = (
            f"\n\n[LIVE WEB SEARCH RESULTS — gathered from {len(results)} independent "
            "source(s), trust-ranked and cross-checked for this query. These results "
            "are NEWER than your training data. Prioritise them over any internal "
            "knowledge.\n\n"
        )
        if snippet_lines:
            context += f"Search Snippets:\n{snippet_lines}\n\n"
        if enriched:
            context += "Detailed Excerpts:\n" + "\n\n".join(
                f"[Source: {r.get('title', 'Unknown')} — {r['url']}]\n{page_text[:1500]}"
                for _, r, page_text in enriched
            ) + "\n"

        if conflict_detected:
            context += (
                f"\n\n[⚠ SOURCE CONFLICT DETECTED: {conflict_detail}. Do NOT merge, "
                "average, or guess which figure is correct. Tell the user this "
                "specific detail could not be confidently verified because trusted "
                "sources disagree, and suggest checking again shortly or consulting "
                "the linked sources directly.]"
            )
        elif confidence == "medium":
            context += (
                "\n\n[VERIFICATION NOTE: Source agreement is moderate. Mention that "
                "this information may change or update soon.]"
            )
        elif confidence == "low":
            context += (
                "\n\n[VERIFICATION NOTE: Confidence in these results is low. Tell the "
                "user this could not be fully verified and present it as provisional, "
                "not confirmed fact.]"
            )

        context += (
            "\nMANDATORY RULES:\n"
            "1. Answer ONLY from these results.\n"
            "2. Do NOT say 'I cannot browse', 'my knowledge cutoff', "
            "or 'I don't have real-time data' — these results ARE your real-time data.\n"
            "3. If a specific detail is absent from the results, say "
            "'I couldn't verify that specific detail' instead of guessing.\n"
            "4. Cite source titles/URLs when stating specific facts.\n"
            "5. If sources conflict on a fact, say so explicitly instead of picking "
            "one arbitrarily.]"
        )

        # Supplement with Wikipedia for background factual depth
        wiki = wikipedia_lookup(query)
        if wiki:
            context += (
                f"\n\n[Wikipedia Background (for factual depth — "
                f"may not reflect the latest developments):\n{wiki}\n]"
            )

    else:
        # Web search returned nothing — fall back to Wikipedia
        wiki = wikipedia_lookup(query)
        if wiki:
            context = (
                f"\n\n[Wikipedia Reference (background knowledge — "
                f"may not reflect recent developments):\n{wiki}\n]"
            )
        elif auto_live:
            context = (
                "\n\n[Live search returned no results for this time-sensitive query. "
                "Tell the user you couldn't verify the current information and suggest "
                "they check a reliable source such as a news site or official page, "
                "instead of guessing.]"
            )
        else:
            context = ""

    # Prediction disclaimer — injected whenever the query asks for a forecast
    prediction_note = (
        "\n\n[PREDICTION MODE ACTIVE: The user is asking for a future prediction, "
        "forecast, or estimate — NOT a verified current fact. You MUST begin your "
        "response with: 'This is a prediction/estimate based on available trends "
        "and data — not a verified fact.' Never present a forecast as a certainty "
        "or as current verified news.]"
    )
    if is_prediction:
        context += prediction_note

    return context


# =============================================================================
# VISION (unchanged — still served via OpenRouter, per feature scope)
# =============================================================================
VISION_MODEL = "openai/gpt-4o"

def call_ai_vision(image_data_url, question, max_tokens=1500):
    """Sends an image + question to a vision-capable model via OpenRouter
    and returns the text reply, or None on failure."""
    if not OPENROUTER_KEY:
        print("Vision call skipped: no OPENROUTER_API_KEY configured")
        return None
    headers = {
        "Authorization": f"Bearer {OPENROUTER_KEY}",
        "Content-Type": "application/json",
        "HTTP-Referer": "https://syra.ai",
        "X-Title": "SYRA AI Platform",
    }
    payload = {
        "model": VISION_MODEL,
        "messages": [
            {"role": "system", "content": SYRA_SYSTEM},
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": question or "Describe this image in detail."},
                    {"type": "image_url", "image_url": {"url": image_data_url}},
                ],
            },
        ],
        "max_tokens": max_tokens,
    }
    try:
        res = requests.post(OPENROUTER_URL, headers=headers, json=payload, timeout=30)
        if res.status_code != 200:
            print("Vision API error:", res.status_code, res.text[:200])
            return None
        return res.json()["choices"][0]["message"]["content"]
    except Exception as e:
        print("Vision call failed:", e)
        return None

# =============================================================================
# SYSTEM PROMPTS & CHAT SYSTEMS
# =============================================================================
SYRA_SYSTEM = """
You are SYRA, an advanced AI assistant created by Surya under the KRISH project.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
IDENTITY RULES (MANDATORY — never break these)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
• Never say you are Google Gemini, ChatGPT, Claude, or any other AI system.
• Never reveal the underlying AI model, provider, or company.
• Never say you are trained by Google, OpenAI, Anthropic, Meta, or anyone else.
• If asked "Who are you?" always reply:
  "I am SYRA, an advanced AI assistant created by Surya under the KRISH project."

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
LIVE SEARCH RULES (when search results are provided)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
• Treat live search results as the absolute current truth — they override training knowledge.
• NEVER say: "I cannot browse the internet."
• NEVER say: "My knowledge cutoff is..."
• NEVER say: "I don't have real-time information."
• Answer ONLY from the provided results when they are present.
• If a specific fact is missing from the results: say "I couldn't verify that specific detail" — never guess.
• Cite the source (title and/or URL) when stating key facts from search results.
• If multiple sources conflict on a specific fact, do NOT merge, average, or guess which one is
  correct. Say plainly that the detail could not be confidently verified, and if useful, mention
  the differing values along with their sources so the user can judge for themselves.
• If the context includes a "[⚠ SOURCE CONFLICT DETECTED...]" or "[VERIFICATION NOTE...]" tag,
  follow its instruction exactly before answering — these reflect an automated cross-check of the
  sources and take priority over your own judgment of how reliable the results look.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
FACTUAL ACCURACY RULES (MANDATORY — zero hallucination)
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
• Never fabricate facts, statistics, names, dates, prices, scores, or quotes.
• If you are not confident in a fact and no search/Wikipedia context confirms it:
  say "I couldn't verify that information." — do not guess or invent.
• It is always better to admit uncertainty than to state an unverified claim as fact.
• Never combine two known facts to produce a third unverified "fact."
• For numbers (prices, scores, percentages, dates): state only what the search results confirm.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
REASONING RULES
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
• Think step-by-step before drawing conclusions on complex questions.
• Verify your reasoning: before stating a conclusion, confirm it is supported by the evidence in context.
• Address multi-part questions systematically — don't skip parts.
• State assumptions explicitly when you make them.
• Prefer the simplest explanation that fits all available evidence.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
SELF-CHECK BEFORE EVERY RESPONSE
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Before sending your reply, internally verify:
✓ Could this answer be outdated? → If yes and no live search data is provided, flag it.
✓ Am I stating something I cannot verify? → If yes, add a disclaimer.
✓ Does my answer contradict the provided search results? → If yes, correct yourself.
✓ Am I making an assumption not in the evidence? → If yes, label it as an assumption.
✓ Have I fabricated any name, number, or quote? → If yes, remove it.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
PREDICTION / FORECAST RULES
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
• When asked for a prediction, forecast, or future estimate, ALWAYS begin with:
  "This is a prediction/estimate based on available trends — not a verified fact."
• Never present a forecast as a certainty.
• Clearly separate: historical facts → current verifiable data → future predictions.
• Base predictions on observable trends and stated reasoning, not speculation.

━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
LANGUAGE & STYLE
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
• Be clear, accurate, and concise.
• Use structured formatting (headings, bullet points) when it improves clarity.
• Match the user's language — respond in Tamil if they write in Tamil or Tanglish.
• Maintain a helpful, knowledgeable, and trustworthy tone at all times.
"""

AGRI_SYSTEM = """
You are a master agricultural consultant and crop science expert.
Provide clear, practical, and localized advice for farming, pest control, crop management, irrigation, and yield improvement.
Structure advice logically:
1. Crop/Soil analysis
2. Seed Selection
3. Fertilizer schedule (Urea, DAP, Potash)
4. Watering/Irrigation plan
5. Disease & Pest management
6. Estimated Market Prices & Harvesting
"""

CODE_SYSTEM = """
You are an expert software developer and architect.
Provide clean, bug-free, fully modular, and commented code.
Ensure proper syntax highlighting.
Explain the architecture, libraries used, and how to run or deploy the code.
"""

WEB_BUILDER_SYSTEM = """
You are a senior frontend developer and UX designer.

You do not explain, discuss, or list technologies. You do not write tutorials
or give advice about how someone else could build the site. You ONLY output
complete, production-ready project files.

OUTPUT FORMAT (MANDATORY — do not deviate):
For every file you generate, use exactly this format, with no text before
the first file and no text after the last file:

=== FILE: index.html ===
<full file content>
=== FILE: style.css ===
<full file content>
=== FILE: script.js ===
<full file content>

Rules:
- Always generate at minimum index.html, style.css, and script.js.
- Add README.md or package.json only if the project genuinely needs them.
- Never wrap file contents in markdown code fences (no ``` characters).
- Never include commentary, explanations, or setup instructions outside the file blocks.
- Build modern, fully responsive, visually stunning interfaces using
  glassmorphism, fluid animations, and clean typography.
- Every file must be complete and immediately runnable — no placeholders
  like "add your code here".

STRICT FILE SEPARATION (MANDATORY):
- index.html must contain ONLY HTML markup. No inline <style> blocks, no
  style="..." attributes, no <script> blocks, no inline onclick="..." or
  other on* handlers. Link the other two files with
  <link rel="stylesheet" href="style.css"> and
  <script src="script.js"></script>.
- style.css must contain ONLY CSS. No HTML tags, no JavaScript.
- script.js must contain ONLY JavaScript. No HTML tags, no CSS rules.
- Every tag opened must be properly closed; output must be valid, well-formed
  HTML5, valid CSS3, and valid modern JavaScript (ES6+) with no syntax errors.
- Use semantic HTML5 elements (header, nav, main, section, article, footer),
  proper heading hierarchy, alt text on images, and ARIA attributes where
  appropriate for accessibility and SEO.
- Include meta viewport and meta description tags for responsiveness and SEO.
- Respect prefers-color-scheme so the site works well in both dark and light
  mode.
- All navigation links, buttons, forms, and interactive cards must be
  functional, not decorative placeholders.
"""

PROJECT_SYSTEM = """
You are an expert full-stack developer.

You do not explain, discuss, or list technologies. You do not write tutorials
or give advice about how someone else could build the project. You ONLY
output complete, production-ready project files.

OUTPUT FORMAT (MANDATORY — do not deviate):
For every file you generate, use exactly this format, with no text before
the first file and no text after the last file:

=== FILE: path/to/file.ext ===
<full file content>

Rules:
- Include all configuration files (package.json, requirements.txt, Dockerfile, etc.) as separate file blocks.
- Include a README.md file block with clear setup, run, and test commands.
- Never wrap file contents in markdown code fences (no ``` characters).
- Never include commentary, explanations, or setup instructions outside the file blocks.
- Every file must be complete and immediately runnable — no placeholders
  like "add your code here".

The README.md file block must cover, in this order:
1. Project overview and architecture (backend, frontend, database)
2. Folder structure
3. Installation guide and dependencies
4. Run instructions (how to start backend/frontend/database locally)
5. API documentation (routes, methods, request/response shape) if the
   project has a backend
6. Future improvements
Never fabricate a library, API, or command that doesn't actually exist —
only reference real, verifiable tools and packages.
"""

# WEBSITE BUILDER WORKSPACE — structured JSON output contract.
# This is intentionally a SEPARATE prompt from WEB_BUILDER_SYSTEM above:
# WEB_BUILDER_SYSTEM stays exactly as-is (still used by the existing
# chat-based /api/website/generate + === FILE: === flow — untouched).
# This one feeds the new dedicated Website Builder workspace panel,
# which needs clean {html, css, javascript} JSON instead of file blocks.
WEBSITE_BUILDER_JSON_SYSTEM = """
You are the SYRA AI Website Builder — a senior frontend engineer and UX designer.

The user will describe a website. Respond with ONE JSON object and NOTHING
else — no markdown fences, no ``` characters, no commentary before or after.

The JSON object must have exactly these three keys:
{
  "html": "...",
  "css": "...",
  "javascript": "..."
}

Rules (MANDATORY):
- Output ONLY the raw JSON object. Never wrap it in code fences.
- All values must be valid JSON strings (escape quotes/newlines properly).
- "html" is always a complete, production-quality, fully responsive page.
- Leave "css" as an empty string ONLY if the page truly needs zero extra
  styling beyond what's already inline in the HTML.
- Leave "javascript" as an empty string ONLY if the page needs zero
  interactivity.
- Never reference external file paths for the CSS/JS you generate — the
  content itself must be the actual stylesheet/script text.
- Build modern, polished, visually stunning interfaces using
  glassmorphism, fluid animation, and clean typography, matching a
  premium SaaS product aesthetic.
- Never include placeholders like "add your code here" — everything must
  be complete and immediately usable.
- The "html" string must contain ONLY markup: no <style> blocks, no
  style="..." attributes, and no <script> blocks or inline on* handlers.
  All CSS goes in "css" and all JavaScript goes in "javascript" — never
  mix languages across the three keys.
- Output must be valid, well-formed HTML5 (every tag closed), valid CSS3,
  and valid modern JavaScript (ES6+) with no syntax errors.
- Use semantic HTML5, proper heading hierarchy, alt text, and ARIA
  attributes for accessibility and SEO, plus a meta viewport tag for
  responsiveness.
"""

# =============================================================================
# PERSISTENT MULTI-CONVERSATION CHAT HISTORY (ChatGPT-style)
# =============================================================================
# Storage layout on disk, per user:
#   MEMORY_FOLDER/<user>/_index.json            -> conversation list/metadata
#   MEMORY_FOLDER/<user>/<conversation_id>.json -> that conversation's messages
#
# NOTE: this is per-process, JSON-file storage. It survives app restarts (all
# reads/writes hit disk) but does NOT synchronize across multiple gunicorn
# workers/dynos — for multi-worker production deployment, swap this for Redis
# or a database-backed store. Flagged here rather than silently shipped as if
# it were already multi-worker safe.
_histories = {}          # in-process cache: (user_key, conv_id) -> messages list
DEFAULT_CONV_ID = "default"  # used when a caller doesn't pass conversation_id,
                              # so existing frontend calls keep working unchanged.

# NOTE ON STORAGE: this used to be per-user JSON files on disk
# (MEMORY_FOLDER/<user>/_index.json + <conv_id>.json). It is now backed by
# the `conversations` and `messages` SQLite tables created in
# init_database(). Every function below keeps its original name and
# signature so none of the call sites (routes, syra_master, etc.) needed
# to change — only the storage engine underneath moved.

def _safe_conv_id(conv_id):
    return secure_filename(conv_id) or DEFAULT_CONV_ID

def _now_iso():
    return datetime.now().isoformat()

def ensure_conversation(user_key, conv_id, first_message=None):
    """Makes sure conv_id has a conversations row, creating one (auto-titled
    from the first message) if it doesn't exist yet. Idempotent."""
    conv_id = _safe_conv_id(conv_id)
    conn = get_db()
    try:
        cur = conn.cursor()
        row = cur.execute(
            "SELECT * FROM conversations WHERE user_key=? AND id=?",
            (user_key, conv_id)
        ).fetchone()
        if row:
            return dict(row)
        title = (first_message or "New chat").strip()[:60] or "New chat"
        now = _now_iso()
        cur.execute(
            """INSERT INTO conversations(id, user_key, title, pinned, created_at, updated_at)
               VALUES (?, ?, ?, 0, ?, ?)""",
            (conv_id, user_key, title, now, now)
        )
        conn.commit()
        return {"id": conv_id, "title": title, "pinned": False, "created_at": now, "updated_at": now}
    finally:
        conn.close()

def touch_conversation(user_key, conv_id):
    conv_id = _safe_conv_id(conv_id)
    conn = get_db()
    try:
        cur = conn.cursor()
        now = _now_iso()
        cur.execute(
            "UPDATE conversations SET updated_at=? WHERE user_key=? AND id=?",
            (now, user_key, conv_id)
        )
        if cur.rowcount == 0:
            conn.close()
            ensure_conversation(user_key, conv_id)
            return
        conn.commit()
    finally:
        conn.close()

def list_conversations(user_key):
    conn = get_db()
    try:
        rows = conn.execute(
            "SELECT * FROM conversations WHERE user_key=? ORDER BY pinned DESC, updated_at DESC",
            (user_key,)
        ).fetchall()
        return [
            {
                "id": r["id"],
                "title": r["title"],
                "pinned": bool(r["pinned"]),
                "created_at": r["created_at"],
                "updated_at": r["updated_at"],
            }
            for r in rows
        ]
    finally:
        conn.close()

def rename_conversation(user_key, conv_id, title):
    conv_id = _safe_conv_id(conv_id)
    title = (title or "").strip()[:120]
    conn = get_db()
    try:
        cur = conn.cursor()
        row = cur.execute(
            "SELECT * FROM conversations WHERE user_key=? AND id=?",
            (user_key, conv_id)
        ).fetchone()
        if not row:
            return None
        new_title = title or row["title"]
        cur.execute(
            "UPDATE conversations SET title=? WHERE user_key=? AND id=?",
            (new_title, user_key, conv_id)
        )
        conn.commit()
        entry = dict(row)
        entry["title"] = new_title
        entry["pinned"] = bool(entry["pinned"])
        return entry
    finally:
        conn.close()

def pin_conversation(user_key, conv_id, pinned):
    conv_id = _safe_conv_id(conv_id)
    conn = get_db()
    try:
        cur = conn.cursor()
        row = cur.execute(
            "SELECT * FROM conversations WHERE user_key=? AND id=?",
            (user_key, conv_id)
        ).fetchone()
        if not row:
            return None
        cur.execute(
            "UPDATE conversations SET pinned=? WHERE user_key=? AND id=?",
            (1 if pinned else 0, user_key, conv_id)
        )
        conn.commit()
        entry = dict(row)
        entry["pinned"] = bool(pinned)
        return entry
    finally:
        conn.close()

def delete_conversation(user_key, conv_id):
    conv_id = _safe_conv_id(conv_id)
    conn = get_db()
    try:
        cur = conn.cursor()
        cur.execute("DELETE FROM conversations WHERE user_key=? AND id=?", (user_key, conv_id))
        deleted = cur.rowcount > 0
        cur.execute("DELETE FROM messages WHERE user_key=? AND conv_id=?", (user_key, conv_id))
        conn.commit()
        _histories.pop((user_key, conv_id), None)
        return deleted
    finally:
        conn.close()

def search_conversations(user_key, query):
    query = (query or "").lower().strip()
    if not query:
        return []
    conn = get_db()
    try:
        like = f"%{query}%"
        rows = conn.execute(
            """SELECT DISTINCT c.* FROM conversations c
               LEFT JOIN messages m ON m.user_key = c.user_key AND m.conv_id = c.id
               WHERE c.user_key = ?
                 AND (LOWER(c.title) LIKE ? OR LOWER(m.content) LIKE ?)
               ORDER BY c.pinned DESC, c.updated_at DESC""",
            (user_key, like, like)
        ).fetchall()
        return [
            {
                "id": r["id"],
                "title": r["title"],
                "pinned": bool(r["pinned"]),
                "created_at": r["created_at"],
                "updated_at": r["updated_at"],
            }
            for r in rows
        ]
    finally:
        conn.close()

def get_history(user_key, conv_id=DEFAULT_CONV_ID):
    conv_id = _safe_conv_id(conv_id)
    cache_key = (user_key, conv_id)
    if cache_key in _histories:
        return _histories[cache_key]
    conn = get_db()
    try:
        rows = conn.execute(
            """SELECT role, content FROM messages
               WHERE user_key=? AND conv_id=? ORDER BY id ASC""",
            (user_key, conv_id)
        ).fetchall()
        history = [{"role": r["role"], "content": r["content"]} for r in rows]
        _histories[cache_key] = history
        return history
    finally:
        conn.close()

def trim_history(user_key, conv_id=DEFAULT_CONV_ID):
    conv_id = _safe_conv_id(conv_id)
    cache_key = (user_key, conv_id)
    h = _histories.get(cache_key, [])
    if len(h) > MAX_HISTORY * 2:
        h = h[-(MAX_HISTORY * 2):]
        _histories[cache_key] = h

    # Persist the full in-memory history to SQLite: clear this conversation's
    # rows and re-insert, which keeps behavior identical to the old
    # "overwrite the whole file" approach while storage is now SQLite.
    conn = get_db()
    try:
        cur = conn.cursor()
        cur.execute("DELETE FROM messages WHERE user_key=? AND conv_id=?", (user_key, conv_id))
        now = _now_iso()
        cur.executemany(
            "INSERT INTO messages(user_key, conv_id, role, content, created_at) VALUES (?, ?, ?, ?, ?)",
            [(user_key, conv_id, m.get("role", "user"), m.get("content", ""), now) for m in h]
        )
        conn.commit()
    except Exception as e:
        print("Failed to persist history:", e)
    finally:
        conn.close()
    touch_conversation(user_key, conv_id)

# =============================================================================
# CHAT HANDLER
# =============================================================================
def syra_master(user_text, history, model_key="default", mode="general", web_search_enabled=False):
    intent = detect_intent(user_text)
    category = map_intent_to_category(intent)
    msg = user_text.lower().strip()

    identity = handle_identity_questions(user_text)
    if identity:
        return identity, "identity", None

    if intent == "image":
        clean = msg
        for w in ["generate image of", "create image of", "image of", "generate image", "create image", "image", "photo", "picture", "draw"]:
            clean = clean.replace(w, "")
        clean = clean.strip() or "highly detailed digital illustration"
        try:
            output = replicate.run("black-forest-labs/flux-schnell", input={"prompt": clean})
            image_url = output[0]
            return {"type": "image", "url": image_url, "prompt": clean}, "image", None
        except Exception as e:
            return f"⚠️ Image generation failed: {str(e)}", "error", None

    if mode == "agri" or intent == "agri":
        system_content = AGRI_SYSTEM
        category = "agri"
    elif mode == "code" or intent == "code":
        system_content = CODE_SYSTEM
        category = "coding"
    elif mode == "website" or intent == "website":
        system_content = WEB_BUILDER_SYSTEM
        category = "website"
    elif mode == "project" or intent == "project":
        system_content = PROJECT_SYSTEM
        category = "project"
    else:
        # General Questions — answered directly by the AI
        system_content = SYRA_SYSTEM
        category = "chat"

    # Smart live search grounding — user-triggered via the toggle, OR
    # auto-triggered for time-sensitive queries (sports, news, weather,
    # markets, politics, technology). build_search_context() handles
    # multi-source trust ranking, URL deduplication, Wikipedia fallback,
    # page-content enrichment, and prediction mode tagging in one call.
    auto_live = needs_live_data(user_text)
    is_prediction = detect_prediction_mode(user_text)
    if web_search_enabled or auto_live:
        print("WEB SEARCH ENABLED", "(auto)" if auto_live and not web_search_enabled else "")
    context_block = build_search_context(
        user_text,
        web_search_enabled=web_search_enabled,
        auto_live=auto_live,
        limit=6,
        is_prediction=is_prediction,
    )
    if context_block:
        system_content += context_block
    elif is_prediction:
        # Prediction query that doesn't need live data — still label as prediction
        system_content += (
            "\n\n[PREDICTION MODE: The user is asking for a future prediction or "
            "estimate. Begin your response with: 'This is a prediction based on "
            "available trends and patterns, not a verified fact.' Never present "
            "a forecast as a certainty.]"
        )

    messages = [{"role": "system", "content": system_content + current_date_context()}] + history + [{"role": "user", "content": user_text}]
    messages.insert(1, {
    "role": "system",
    "content": "Never say you are Google Gemini, ChatGPT, Claude or OpenAI. Always identify yourself only as SYRA, created by Surya under the KRISH project."
})
    # Centralized AI routing with category-aware multi-provider fallback
    reply = generate_ai_reply(messages, category=category, max_tokens=4096)

    history.append({"role": "user", "content": user_text})
    history.append({"role": "assistant", "content": reply})
    return reply, intent, None

def resolve_session_id(req):
    """Single source of truth for user bucketing across endpoints.
    Logged-in users always key off their Flask session email so /api/chat
    and /api/chat/stream share the same conversation memory."""
    flask_user_email = session.get("user", {}).get("email")
    return flask_user_email or req.get("session_id") or "anon"

def resolve_conversation_id(req, user_key, first_message=None):
    """Resolves which conversation a request belongs to. Falls back to a
    single DEFAULT_CONV_ID when the frontend doesn't send conversation_id,
    so older frontend builds keep working exactly as before. Auto-creates
    the conversation's index entry (auto-titled from the first message) if
    it doesn't exist yet."""
    conv_id = (req.get("conversation_id") or DEFAULT_CONV_ID).strip() or DEFAULT_CONV_ID
    conv_id = secure_filename(conv_id) or DEFAULT_CONV_ID
    ensure_conversation(user_key, conv_id, first_message=first_message)
    return conv_id

# =============================================================================
# CHAT ENDPOINTS
# =============================================================================
@app.route("/api/chat", methods=["POST"])
def api_chat():
    print("API CHAT CALLED")
    try:
        req = request.get_json() or {}
        user_text = req.get("message", "").strip()
        model_key = req.get("model", "default")
        session_id = resolve_session_id(req)
        mode = req.get("mode", "general")
        web_search_enabled = req.get("web_search", False)

        if not user_text:
            return jsonify(ok=False, message="Empty message"), 400

        conv_id = resolve_conversation_id(req, session_id, first_message=user_text)
        history = get_history(session_id, conv_id)
        result, res_mode, _ = syra_master(user_text, history, model_key, mode, web_search_enabled)
        trim_history(session_id, conv_id)

        if isinstance(result, dict):
            return jsonify(ok=True, mode=res_mode, conversation_id=conv_id, data=result, message="Success")

        return jsonify(ok=True, mode=res_mode, conversation_id=conv_id, message=str(result))
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500

@app.route("/api/chat/stream", methods=["POST"])
def api_chat_stream():
    try:
        print("STREAM API CALLED")

        req = request.get_json() or {}
        user_text = req.get("message", "").strip()

        print("USER:", user_text)

        model_key = req.get("model", "default")
        session_id = resolve_session_id(req)
        mode = req.get("mode", "general")
        web_search_enabled = req.get("web_search", False)

        print("WEB SEARCH:", web_search_enabled)

        if not user_text:
            return jsonify(ok=False, message="Empty message"), 400

        conv_id = resolve_conversation_id(req, session_id, first_message=user_text)

        # Detect image intent — return non-streaming JSON event
        intent = detect_intent(user_text)

        if intent == "image":
            clean = user_text.lower()
            for w in ["generate image of", "create image of", "image of", "generate image", "create image", "image", "photo", "picture", "draw"]:
                clean = clean.replace(w, "")
            clean = clean.strip() or "highly detailed digital illustration"
            image_url, err = generate_image_with_fallback(clean)
            if image_url:
                def image_event():
                    yield f"data: {json.dumps({'conversation_id': conv_id})}\n\n"
                    yield f"data: {json.dumps({'type': 'image', 'url': image_url, 'prompt': clean})}\n\n"
                    yield "data: [DONE]\n\n"
                return Response(stream_with_context(image_event()), content_type="text/event-stream")
            else:
                print("IMAGE GENERATION ERROR:", err)

                def err_event():
                    yield f"data: {json.dumps({'conversation_id': conv_id})}\n\n"
                    yield f"data: {json.dumps({'text': err or FRIENDLY_ERROR})}\n\n"
                    yield "data: [DONE]\n\n"

                return Response(
                    stream_with_context(err_event()),
                    content_type="text/event-stream"
                )

        identity = handle_identity_questions(user_text)
        if identity:
            def identity_event():
                yield f"data: {json.dumps({'conversation_id': conv_id})}\n\n"
                yield f"data: {json.dumps({'text': identity})}\n\n"
                yield "data: [DONE]\n\n"
            return Response(stream_with_context(identity_event()), content_type="text/event-stream")

        system_map = {
            "agri": AGRI_SYSTEM,
            "code": CODE_SYSTEM,
            "website": WEB_BUILDER_SYSTEM,
            "project": PROJECT_SYSTEM,
        }
        system_content = system_map.get(mode, SYRA_SYSTEM)

        # Smart live search grounding — user-triggered via the toggle, OR
        # auto-triggered for time-sensitive queries (sports, news, weather,
        # markets, politics, technology). Shares the same build_search_context()
        # helper used by syra_master() for consistent behaviour across both
        # streaming and non-streaming chat paths.
        auto_live = needs_live_data(user_text)
        is_prediction = detect_prediction_mode(user_text)
        if web_search_enabled or auto_live:
            print("WEB SEARCH ENABLED (stream)", "(auto)" if auto_live and not web_search_enabled else "")
        context_block = build_search_context(
            user_text,
            web_search_enabled=web_search_enabled,
            auto_live=auto_live,
            limit=6,
            is_prediction=is_prediction,
        )
        if context_block:
            system_content += context_block
        elif is_prediction:
            system_content += (
                "\n\n[PREDICTION MODE: The user is asking for a future prediction or "
                "estimate. Begin your response with: 'This is a prediction based on "
                "available trends and patterns, not a verified fact.' Never present "
                "a forecast as a certainty.]"
            )

        system_content += current_date_context()

        history = get_history(session_id, conv_id)
        
        # Determine request category for smart provider routing
        intent = detect_intent(user_text)
        category = map_intent_to_category(intent)
        if mode == "agri":
            category = "agri"
        elif mode == "code":
            category = "coding"
        elif mode == "website":
            category = "website"
        elif mode == "project":
            category = "project"
        
        messages = [{"role": "system", "content": system_content}] + history + [{"role": "user", "content": user_text}]
        messages.insert(1, {
            "role": "system",
            "content": "Never say you are Google Gemini, ChatGPT, Claude or OpenAI. Always identify yourself only as SYRA, created by Surya under the KRISH project."
        })

        def generate():
            yield f"data: {json.dumps({'conversation_id': conv_id})}\n\n"
            full = ""
            # Centralized AI routing with category-aware multi-provider fallback
            for chunk in generate_ai_reply_stream(messages, category=category, max_tokens=4096):
                yield chunk
                if '"text":' in chunk:
                    try:
                        full += json.loads(chunk[6:]).get("text", "")
                    except Exception:
                        pass
            history.append({"role": "user", "content": user_text})
            history.append({"role": "assistant", "content": full})
            trim_history(session_id, conv_id)

        return Response(stream_with_context(generate()), content_type="text/event-stream")

    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500


# =============================================================================
# CONVERSATION HISTORY ENDPOINTS (ChatGPT-style sidebar)
# =============================================================================
@app.route("/api/conversations", methods=["GET"])
def api_conversations_list():
    try:
        user_key = resolve_session_id({"session_id": request.args.get("session_id")})
        return jsonify(ok=True, conversations=list_conversations(user_key))
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, error=str(e)), 500

@app.route("/api/conversations", methods=["POST"])
def api_conversations_create():
    try:
        req = request.get_json() or {}
        user_key = resolve_session_id(req)
        conv_id = str(uuid.uuid4())[:12]
        title = req.get("title") or "New chat"
        entry = ensure_conversation(user_key, conv_id, first_message=title)
        return jsonify(ok=True, conversation=entry)
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, error=str(e)), 500

@app.route("/api/conversations/search", methods=["GET"])
def api_conversations_search():
    try:
        user_key = resolve_session_id({"session_id": request.args.get("session_id")})
        query = request.args.get("q", "")
        return jsonify(ok=True, conversations=search_conversations(user_key, query))
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, error=str(e)), 500

@app.route("/api/conversations/<conv_id>", methods=["GET"])
def api_conversations_get(conv_id):
    try:
        user_key = resolve_session_id({"session_id": request.args.get("session_id")})
        safe_id = secure_filename(conv_id)
        if not any(e["id"] == safe_id for e in list_conversations(user_key)):
            return jsonify(ok=False, error="Conversation not found"), 404
        return jsonify(ok=True, conversation_id=safe_id, messages=get_history(user_key, safe_id))
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, error=str(e)), 500

@app.route("/api/conversations/<conv_id>", methods=["PATCH"])
def api_conversations_update(conv_id):
    try:
        req = request.get_json() or {}
        user_key = resolve_session_id(req)
        safe_id = secure_filename(conv_id)
        entry = None
        if "title" in req:
            entry = rename_conversation(user_key, safe_id, req.get("title", ""))
        if "pinned" in req:
            entry = pin_conversation(user_key, safe_id, req.get("pinned"))
        if entry is None:
            return jsonify(ok=False, error="Conversation not found"), 404
        return jsonify(ok=True, conversation=entry)
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, error=str(e)), 500

@app.route("/api/conversations/<conv_id>", methods=["DELETE"])
def api_conversations_delete(conv_id):
    try:
        user_key = resolve_session_id({"session_id": request.args.get("session_id")})
        safe_id = secure_filename(conv_id)
        deleted = delete_conversation(user_key, safe_id)
        if not deleted:
            return jsonify(ok=False, error="Conversation not found"), 404
        return jsonify(ok=True, deleted=safe_id)
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, error=str(e)), 500

# =============================================================================
# IMAGE GENERATION AI
# =============================================================================
# =============================================================================
# IMAGE GENERATION — Cloudflare Workers AI (primary) → Replicate (fallback)
# =============================================================================
class CloudflareQuotaExceeded(Exception):
    """Raised when Cloudflare's daily free image-generation quota is hit."""
    pass

def _cloudflare_generate_image(prompt, width=1024, height=1024):
    """Calls Cloudflare Workers AI for text-to-image generation. Returns a
    data: URL (base64 PNG/JPEG) on success. Returns None if Cloudflare
    isn't configured or the call fails for a non-quota reason. Raises
    CloudflareQuotaExceeded specifically when the daily free quota is hit,
    so the caller can show the correct message instead of silently falling
    back to Replicate in that specific case."""
    if not (CLOUDFLARE_ACCOUNT_ID and CLOUDFLARE_API_TOKEN):
        return None
    try:
        url = (
            f"https://api.cloudflare.com/client/v4/accounts/"
            f"{CLOUDFLARE_ACCOUNT_ID}/ai/run/{CLOUDFLARE_IMAGE_MODEL}"
        )
        headers = {"Authorization": f"Bearer {CLOUDFLARE_API_TOKEN}"}
        payload = {"prompt": prompt, "width": width, "height": height}
        res = requests.post(url, headers=headers, json=payload, timeout=45)

        if res.status_code == 429:
            raise CloudflareQuotaExceeded()

        if res.status_code != 200:
            print(f"[image] Cloudflare error: {res.status_code} {res.text[:200]}")
            return None

        content_type = res.headers.get("Content-Type", "")
        if content_type.startswith("image/"):
            b64 = base64.b64encode(res.content).decode("utf-8")
            return f"data:{content_type};base64,{b64}"

        # Some Workers AI responses come back as JSON with a quota/error message.
        try:
            data = res.json()
        except Exception:
            return None
        errors = data.get("errors") or []
        if any("quota" in str(e).lower() or "limit" in str(e).lower() for e in errors):
            raise CloudflareQuotaExceeded()
        result = data.get("result") or {}
        img_b64 = result.get("image")
        if img_b64:
            return f"data:image/png;base64,{img_b64}"
        return None
    except CloudflareQuotaExceeded:
        raise
    except requests.Timeout:
        print("[image] Cloudflare timeout")
        return None
    except Exception as e:
        print(f"[image] Cloudflare error: {e}")
        return None

def generate_image_with_fallback(prompt, width=1024, height=1024):
    """Primary/fallback image generation used across the app: Cloudflare
    Workers AI first, Replicate (FLUX) second. Returns (url_or_data_url,
    error_message_or_None). On the specific case of Cloudflare's daily free
    quota being exhausted, still tries Replicate before giving up, but if
    Replicate ALSO fails, returns the friendly daily-limit message rather
    than a raw Replicate error, since quota exhaustion was the root cause."""
    quota_hit = False
    try:
        cf_result = _cloudflare_generate_image(prompt, width, height)
        if cf_result:
            return cf_result, None
    except CloudflareQuotaExceeded:
        quota_hit = True
        print("[image] Cloudflare daily quota exhausted — falling back to Replicate")

    try:
        output = replicate.run(
            "black-forest-labs/flux-schnell",
            input={"prompt": prompt, "width": width, "height": height}
        )
        return output[0], None
    except Exception as e:
        if quota_hit:
            return None, "⚠️ Daily free image generation limit reached.\nPlease try again tomorrow."
        return None, f"⚠️ Image generation failed: {str(e)}"

@app.route("/api/image/generate", methods=["POST"])
def api_image_generate():
    req = request.get_json() or {}
    prompt = req.get("prompt", "").strip()
    width = req.get("width", 1024)
    height = req.get("height", 1024)
    if not prompt:
        return jsonify(ok=False, error="prompt is required"), 400
    url, err = generate_image_with_fallback(prompt, width, height)
    if err:
        return jsonify(ok=False, error=err), 500
    return jsonify(ok=True, url=url, prompt=prompt)

@app.route("/api/image/variations", methods=["POST"])
def api_image_variations():
    req = request.get_json() or {}
    prompt = req.get("prompt", "").strip()
    count = min(int(req.get("count", 3)), 4)
    if not prompt:
        return jsonify(ok=False, error="prompt is required"), 400
    results = []
    for i in range(count):
        try:
            out = replicate.run("black-forest-labs/flux-schnell", input={"prompt": prompt, "seed": i * 11})
            results.append(out[0])
        except Exception as e:
            results.append(f"Error: {str(e)}")
    return jsonify(ok=True, urls=results, prompt=prompt)

@app.route("/api/image/upscale", methods=["POST"])
def api_image_upscale():
    req = request.get_json() or {}
    image_url = req.get("url", "").strip()
    if not image_url:
        return jsonify(ok=False, error="image url is required"), 400
    try:
        out = replicate.run(
            "nightmareai/real-esrgan:42fed1c4974146d4d2414e2be2c5277c7fcf05fcc3a73abf41610695738c1d7b",
            input={"image": image_url, "scale": 4}
        )
        return jsonify(ok=True, url=str(out))
    except Exception as e:
        return jsonify(ok=False, error=str(e)), 500

@app.route("/api/image/remove-background", methods=["POST"])
def api_image_remove_bg():
    req = request.get_json() or {}
    image_url = req.get("url", "").strip()
    if not image_url:
        return jsonify(ok=False, error="image url is required"), 400
    try:
        out = replicate.run(
            "cjwbw/rembg:fb8af171cfa1616ddcf1242c093f9c46bcada5ad4cf6f2fbe8b81b330ec5c003",
            input={"image": image_url}
        )
        return jsonify(ok=True, url=str(out))
    except Exception as e:
        return jsonify(ok=False, error=str(e)), 500

# =============================================================================
# FILE UPLOAD — DATASET (was missing from backend, called by frontend)
# =============================================================================
@app.route("/api/upload_dataset", methods=["POST"])
def api_upload_dataset():
  try:
    if "file" not in request.files:
        return jsonify(ok=False, success=False, error="No file part in request"), 400
    file = request.files["file"]
    if not file or file.filename == "":
        return jsonify(ok=False, success=False, error="No file selected"), 400
    if not allowed_file(file.filename):
        return jsonify(ok=False, success=False, error="File type not allowed"), 400

    filename = secure_filename(file.filename)
    save_path = os.path.join(UPLOAD_FOLDER, filename)
    file.save(save_path)

    ext = filename.rsplit(".", 1)[-1].lower()

    # Images are routed to Vision instead of the text-preview pipeline below —
    # the frontend calls /api/vision separately with this same file, so we
    # just report the category here and skip trying to "preview" binary
    # image bytes as text.
    if ext in IMAGE_EXTENSIONS:
        # Auto-route image uploads straight to Vision analysis — no manual
        # "analyze this image" step required. Vision failures never fail the
        # upload itself; they just leave vision_analysis empty so existing
        # callers relying only on the original fields keep working.
        vision_analysis = None
        try:
            with open(save_path, "rb") as f:
                raw = f.read()
            mime = "image/jpeg" if ext == "jpg" else f"image/{ext}"
            image_data_url = f"data:{mime};base64,{base64.b64encode(raw).decode('utf-8')}"
            vision_analysis = call_ai_vision(image_data_url, "Describe this image in detail.")
        except Exception as e:
            print("Auto vision routing failed:", e)

        return jsonify(
            ok=True,
            success=True,
            filename=filename,
            size=os.path.getsize(save_path),
            category="image",
            preview="[Image uploaded — ready for Vision analysis]",
            vision_analysis=vision_analysis
        )

    category = "code" if ext in CODE_EXTENSIONS else "document"

    # Extract text summary for non-binary files
    preview = ""
    try:
        if ext in {"txt", "md", "csv", "json", "xml"} or ext in CODE_EXTENSIONS:
            with open(save_path, "r", encoding="utf-8", errors="ignore") as f:
                preview = f.read(2000)
        elif ext == "pdf":
            try:
                import PyPDF2
                with open(save_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    preview = " ".join(page.extract_text() or "" for page in reader.pages[:3])[:2000]
            except Exception:
                preview = "[PDF uploaded — text extraction unavailable]"
        elif ext in {"xlsx", "xls"}:
            try:
                df = pd.read_excel(save_path, nrows=20)
                preview = df.to_string()[:2000]
            except Exception:
                preview = "[Spreadsheet uploaded]"
        elif ext == "docx":
            try:
                import docx
                doc = docx.Document(save_path)
                preview = "\n".join(p.text for p in doc.paragraphs[:30])[:2000]
            except Exception:
                preview = "[Word document uploaded]"
        elif ext == "pptx":
            try:
                from pptx import Presentation
                prs = Presentation(save_path)
                slides_text = []
                for i, slide in enumerate(prs.slides[:20] if hasattr(prs.slides, "__getitem__") else list(prs.slides)[:20]):
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text:
                            slides_text.append(shape.text)
                preview = "\n".join(slides_text)[:2000]
            except Exception:
                preview = "[PowerPoint uploaded — text extraction unavailable]"
        elif ext == "zip":
            try:
                with zipfile.ZipFile(save_path, "r") as zf:
                    names = zf.namelist()[:50]
                    preview = "Archive contents:\n" + "\n".join(names)
            except Exception:
                preview = "[ZIP uploaded — could not list contents]"
        else:
            preview = f"[{ext.upper()} file uploaded successfully]"
    except Exception as e:
        preview = f"[File uploaded — preview unavailable: {str(e)}]"

    return jsonify(
        ok=True,
        success=True,
        filename=filename,
        size=os.path.getsize(save_path),
        category=category,
        preview=preview[:500] if preview else ""
    )
  except Exception as e:
    traceback.print_exc()
    return jsonify(ok=False, success=False, error=str(e)), 500

# =============================================================================
# VISION AI
# =============================================================================
@app.route("/api/vision", methods=["POST"])
def api_vision():
    """Accepts either a multipart image upload (field 'image') or a JSON
    body with an 'image_url'/'image_data_url', plus an optional 'question'."""
    try:
        question = ""
        image_data_url = None

        if "image" in request.files:
            file = request.files["image"]
            if not file or file.filename == "":
                return jsonify(ok=False, error="No image selected"), 400
            ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
            if ext not in IMAGE_EXTENSIONS:
                return jsonify(ok=False, error="Unsupported image type"), 400
            raw = file.read()
            if len(raw) > 20 * 1024 * 1024:
                return jsonify(ok=False, error="Image too large (max 20MB)"), 400
            mime = "image/jpeg" if ext == "jpg" else f"image/{ext}"
            image_data_url = f"data:{mime};base64,{base64.b64encode(raw).decode('utf-8')}"
            question = request.form.get("question", "")
        else:
            req = request.get_json(silent=True) or {}
            image_data_url = req.get("image_url") or req.get("image_data_url")
            question = req.get("question", "")

        if not image_data_url:
            return jsonify(ok=False, error="image is required"), 400

        reply = call_ai_vision(image_data_url, question or "Describe this image in detail.")
        if not reply:
            return jsonify(ok=False, error="Vision analysis failed — check OPENROUTER_API_KEY is configured for a vision-capable model."), 500

        return jsonify(ok=True, message=reply)
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, error=str(e)), 500

# =============================================================================
# WEBSITE & PROJECT BUILDERS
# =============================================================================
def parse_project_files(content):
    files = {}
    # Defensive cleanup: some models wrap the ENTIRE reply in a single
    # markdown fence around the "=== FILE: ..." blocks despite being told
    # not to. Strip that outer fence before the existing per-file parsing
    # below (which already strips fences around individual files).
    content = re.sub(r"^\s*```[a-zA-Z0-9]*\n", "", content)
    content = re.sub(r"\n```\s*$", "", content)
    pattern = r"===\s*FILE:\s*(.+?)\s*===\n(.*?)(?====\s*FILE:|\Z)"
    matches = re.findall(pattern, content, re.DOTALL | re.IGNORECASE)
    for path, code in matches:
        path = path.strip()
        code = code.strip()
        # Defensive cleanup: strip a leading/trailing markdown code fence
        # in case the model adds one despite being told not to.
        code = re.sub(r"^```[a-zA-Z0-9]*\n", "", code)
        code = re.sub(r"\n```$", "", code)
        files[path] = code.strip()
    if not files:
        # Model ignored the file-block format entirely — fall back to a
        # single plain-text file rather than silently discarding the reply.
        # IMPORTANT: this raw content is untrusted model output and must
        # never be labeled/served as index.html, or the Project Explorer's
        # HTML preview would execute it as live markup.
        files["output.txt"] = content
    return files

def make_zip(project_id, files):
    proj_dir = os.path.join(PROJECTS_FOLDER, project_id)
    os.makedirs(proj_dir, exist_ok=True)
    zip_path = os.path.join(PROJECTS_FOLDER, f"{project_id}.zip")
    with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as zf:
        for rel, data in files.items():
            full = os.path.join(proj_dir, rel)
            os.makedirs(os.path.dirname(full), exist_ok=True)
            with open(full, "w", encoding="utf-8") as f:
                f.write(data)
            zf.write(full, rel)
    return zip_path

@app.route("/api/website/generate", methods=["POST"])
def api_website_generate():
    req = request.get_json() or {}
    prompt = req.get("prompt", "").strip()
    stack = req.get("stack", "html")
    if not prompt:
        return jsonify(ok=False, error="prompt is required"), 400

    prompt_pkg = f"Create a gorgeous single-page website with the topic: {prompt}. Stack: {stack}."

    # Retry loop: the model occasionally explains how to build the site
    # instead of emitting the "=== FILE: ..." blocks. parse_project_files()
    # falls back to a single output.txt in that case — treat that as an
    # invalid generation and retry (existing parser is untouched).
    max_attempts = 3
    reply = None
    files = {}
    for attempt in range(max_attempts):
        candidate = generate_ai_reply([{"role": "system", "content": WEB_BUILDER_SYSTEM}, {"role": "user", "content": prompt_pkg}], category="website", max_tokens=6000)
        if not candidate:
            continue
        candidate_files = parse_project_files(candidate)
        if list(candidate_files.keys()) != ["output.txt"]:
            reply, files = candidate, candidate_files
            break
        # Keep the last attempt as a fallback in case every retry fails.
        reply, files = candidate, candidate_files

    if not reply:
        return jsonify(ok=False, error="AI generation failed"), 500

    # Safety net: if the model mixed HTML/CSS/JS across the standard
    # index.html / style.css / script.js trio, cleanly separate them.
    if "index.html" in files:
        clean_html, clean_css, clean_js = separate_html_css_js(
            files.get("index.html", ""), files.get("style.css", ""), files.get("script.js", "")
        )
        files["index.html"] = clean_html
        if clean_css.strip():
            files["style.css"] = clean_css
        if clean_js.strip():
            files["script.js"] = clean_js

    proj_id = str(uuid.uuid4())[:8]
    make_zip(proj_id, files)

    return jsonify(ok=True, project_id=proj_id, files=list(files.keys()), code=reply, download=f"/api/project/download/{proj_id}")

@app.route("/api/project/generate", methods=["POST"])
def api_project_generate():
    req = request.get_json() or {}
    prompt = req.get("prompt", "").strip()
    tech = req.get("tech", "React")
    if not prompt:
        return jsonify(ok=False, error="prompt is required"), 400

    prompt_pkg = f"Generate a full operational software system. Techstack: {tech}. Purpose: {prompt}."
    reply = generate_ai_reply([{"role": "system", "content": PROJECT_SYSTEM}, {"role": "user", "content": prompt_pkg}], category="project", max_tokens=8000)

    if not reply:
        return jsonify(ok=False, error="AI generation failed"), 500

    proj_id = str(uuid.uuid4())[:8]
    files = parse_project_files(reply)
    make_zip(proj_id, files)

    return jsonify(ok=True, project_id=proj_id, files=list(files.keys()), code=reply, download=f"/api/project/download/{proj_id}")

@app.route("/api/project/download/<project_id>")
def api_project_download(project_id):
    pid = secure_filename(project_id)
    zip_path = os.path.join(PROJECTS_FOLDER, f"{pid}.zip")
    if not os.path.exists(zip_path):
        return jsonify(ok=False, error="Project ZIP not found"), 404
    return send_from_directory(PROJECTS_FOLDER, f"{pid}.zip", as_attachment=True)

# =============================================================================
# WEBSITE BUILDER WORKSPACE (new, additive) — dedicated tabbed HTML/CSS/JS
# panel in the frontend. Fully independent of the chat-based website
# generator above: separate system prompt, separate route, separate
# response shape. Does not touch conversations/messages storage, auth,
# or any other existing feature.
# =============================================================================
def separate_html_css_js(html, css, js):
    """Defensive safety net for the Website Builder: even with the stricter
    prompt rules, a model can still occasionally leave a <style> or
    <script> block inside the HTML. This pulls any such blocks out of the
    HTML and merges their content into css/js instead of discarding them,
    so the three files the user gets are always cleanly separated.
    Never used to reject a generation — only to clean it up."""
    html = html or ""
    css = css or ""
    js = js or ""
    try:
        soup = BeautifulSoup(html, "html.parser")

        extracted_css = []
        for style_tag in soup.find_all("style"):
            if style_tag.string:
                extracted_css.append(style_tag.string)
            style_tag.decompose()

        extracted_js = []
        for script_tag in soup.find_all("script"):
            if not script_tag.get("src") and script_tag.string:
                extracted_js.append(script_tag.string)
            script_tag.decompose()

        # Pull inline style="..." attributes out into CSS via generated
        # utility classes only when present; otherwise leave HTML untouched.
        for i, tag in enumerate(soup.find_all(style=True)):
            inline_rule = tag.get("style")
            if not inline_rule:
                continue
            util_class = f"syra-inline-{i}"
            existing_classes = tag.get("class", [])
            tag["class"] = existing_classes + [util_class]
            del tag["style"]
            extracted_css.append(f".{util_class} {{ {inline_rule} }}")

        cleaned_html = str(soup)

        if extracted_css:
            css = (css + "\n\n" if css.strip() else "") + "\n".join(extracted_css)
        if extracted_js:
            js = (js + "\n\n" if js.strip() else "") + "\n".join(extracted_js)

        return cleaned_html, css, js
    except Exception as e:
        print("separate_html_css_js failed, returning original files:", e)
        return html, css, js


def parse_website_builder_json(raw_text):
    """Defensively parses the AI's structured Website Builder reply into
    {html, css, javascript} strings. Handles the model wrapping the JSON
    in a markdown fence, adding stray text around it, or (worst case)
    ignoring the JSON contract entirely — always returns all three keys
    so the frontend never has to special-case a malformed response."""
    if not raw_text:
        return {"html": "", "css": "", "javascript": ""}

    text = raw_text.strip()
    # Strip an outer ```json ... ``` / ``` ... ``` fence if the model added one.
    text = re.sub(r"^```[a-zA-Z0-9]*\s*\n", "", text)
    text = re.sub(r"\n```\s*$", "", text)
    text = text.strip()

    data = None
    try:
        data = json.loads(text)
    except Exception:
        # Model added commentary around the JSON — grab the first {...} block.
        match = re.search(r"\{.*\}", text, re.DOTALL)
        if match:
            try:
                data = json.loads(match.group(0))
            except Exception:
                data = None

    if not isinstance(data, dict):
        # Model ignored the JSON contract entirely — don't discard the
        # reply, treat it as the HTML so the user still gets something.
        return {"html": text, "css": "", "javascript": ""}

    return {
        "html": str(data.get("html", "") or ""),
        "css": str(data.get("css", "") or ""),
        "javascript": str(data.get("javascript") or data.get("js") or ""),
    }

@app.route("/api/website-builder/generate", methods=["POST"])
def api_website_builder_generate():
    """Dedicated endpoint for the AI Website Builder workspace panel.
    Always returns structured {html, css, javascript} JSON for the
    tabbed editor UI, instead of the === FILE: === blocks used by the
    older chat-based /api/website/generate flow (left untouched)."""
    try:
        req = request.get_json() or {}
        prompt = (req.get("prompt") or "").strip()
        if not prompt:
            return jsonify(ok=False, message="Please describe the website you want."), 400

        user_prompt = f"Build this website: {prompt}"
        messages = [
            {"role": "system", "content": WEBSITE_BUILDER_JSON_SYSTEM + current_date_context()},
            {"role": "user", "content": user_prompt},
        ]

        max_attempts = 3
        raw = None
        for attempt in range(max_attempts):
            candidate = generate_ai_reply(messages, category="website", max_tokens=6000)
            raw = candidate
            if candidate and candidate.strip() != FRIENDLY_ERROR.strip():
                break

        if not raw or raw.strip() == FRIENDLY_ERROR.strip():
            return jsonify(ok=False, message=FRIENDLY_ERROR), 503

        result = parse_website_builder_json(raw)

        if not result["html"].strip():
            return jsonify(ok=False, message="The AI didn't return any HTML. Please try again."), 500

        # Safety net: strip out any HTML/CSS/JS the model still mixed
        # together despite the prompt rules, merging it into the right file.
        clean_html, clean_css, clean_js = separate_html_css_js(
            result["html"], result["css"], result["javascript"]
        )

        return jsonify(ok=True, html=clean_html, css=clean_css, javascript=clean_js)
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500

# =============================================================================
# FILE EXPLORER / FILE SYSTEM
# =============================================================================
@app.route("/api/fs/list", methods=["GET"])
def api_fs_list():
    items = []
    for f in os.listdir(UPLOAD_FOLDER):
        full_path = os.path.join(UPLOAD_FOLDER, f)
        items.append({
            "name": f,
            "type": "directory" if os.path.isdir(full_path) else "file",
            "size": os.path.getsize(full_path) if os.path.isfile(full_path) else 0
        })
    return jsonify(ok=True, items=items)

@app.route("/api/fs/read", methods=["POST"])
def api_fs_read():
    req = request.get_json() or {}
    name = secure_filename(req.get("name", ""))
    path = os.path.join(UPLOAD_FOLDER, name)
    if not os.path.exists(path):
        return jsonify(ok=False, error="File not found"), 404
    try:
        with open(path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        return jsonify(ok=True, name=name, content=content)
    except Exception as e:
        return jsonify(ok=False, error=str(e)), 500

@app.route("/api/fs/write", methods=["POST"])
def api_fs_write():
    req = request.get_json() or {}
    name = secure_filename(req.get("name", ""))
    content = req.get("content", "")
    if not name:
        return jsonify(ok=False, error="Name is required"), 400
    try:
        path = os.path.join(UPLOAD_FOLDER, name)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        return jsonify(ok=True, name=name)
    except Exception as e:
        return jsonify(ok=False, error=str(e)), 500

@app.route("/api/fs/delete", methods=["POST"])
def api_fs_delete():
    req = request.get_json() or {}
    name = secure_filename(req.get("name", ""))
    path = os.path.join(UPLOAD_FOLDER, name)
    if not os.path.exists(path):
        return jsonify(ok=False, error="File not found"), 404
    try:
        if os.path.isdir(path):
            shutil.rmtree(path)
        else:
            os.remove(path)
        return jsonify(ok=True, deleted=name)
    except Exception as e:
        return jsonify(ok=False, error=str(e)), 500

# =============================================================================
# VOICE & AUDIO AI
# =============================================================================
@app.route("/api/voice/stt", methods=["POST"])
def api_voice_stt():

    if "audio" not in request.files:
        return jsonify(ok=False, error="No audio file"), 400

    audio_file = request.files["audio"]

    filename = secure_filename(audio_file.filename or "audio.wav")
    path = os.path.join(VOICES_FOLDER, filename)

    audio_file.save(path)

    try:

        transcript = cloudflare_speech_to_text(path)

        return jsonify(
            ok=True,
            transcript=transcript
        )

    except Exception as e:

        return jsonify(
            ok=False,
            error=str(e)
        ), 500

# =============================================================================
# WEB SEARCH
# =============================================================================
@app.route("/api/search", methods=["POST"])
def api_search():
    req = request.get_json() or {}
    query = req.get("query", "").strip()
    if not query:
        return jsonify(ok=False, error="query is required"), 400
    results = web_search(query)
    return jsonify(ok=True, results=results)

# =============================================================================
# AUTH — OWN-BACKEND AUTHENTICATION (Firebase/Supabase fully removed)
# =============================================================================
# Users are stored as a single JSON file of {email: {name, password_hash,
# created_at, reset_token, reset_expires}}. This matches the JSON-file
# persistence pattern already used for chat history/memory elsewhere in this
# file. For multi-worker production deployment, swap this for a real
# database — flagged here rather than silently shipped as if it already were
# one.


def validate_email(email):
    if not email:
        return False

    return bool(
        re.match(
            r"^[^@\s]+@[^@\s]+\.[^@\s]+$",
            email.strip()
        )
    )


def start_user_session(email, name, remember=False):
    session.clear()
    session["user"] = {"email": email, "name": name}
    session["csrf_token"] = secrets.token_hex(32)
    session.permanent = bool(remember)

# =============================================================================
# CSRF PROTECTION (double-submit token)
# =============================================================================
# A random token is minted into the (HttpOnly) session on the first request
# and handed to the frontend via GET /api/csrf-token. Every state-changing
# request (POST/PUT/PATCH/DELETE) to /api/* must echo that same token back in
# the X-CSRF-Token header. An attacker's cross-site page can trigger a request
# with the user's cookies attached, but it cannot read /api/csrf-token's
# response (blocked by CORS/same-origin policy) or the HttpOnly cookie, so it
# can never produce a matching header.
CSRF_EXEMPT_PATHS = set()  # nothing exempt — even login/register require the
                            # pre-session token minted on first page load.

@app.before_request
def _ensure_csrf_token():
    if "csrf_token" not in session:
        session["csrf_token"] = secrets.token_hex(32)

@app.before_request
def _enforce_csrf():
    if request.method not in ("POST", "PUT", "PATCH", "DELETE"):
        return None
    if not request.path.startswith("/api/"):
        return None
    if request.path in CSRF_EXEMPT_PATHS:
        return None
    sent = request.headers.get("X-CSRF-Token", "")
    expected = session.get("csrf_token", "")
    if not sent or not expected or not hmac.compare_digest(sent, expected):
        return jsonify(ok=False, success=False, error="Invalid or missing CSRF token"), 403
    return None

@app.route("/api/csrf-token", methods=["GET"])
def api_csrf_token():
    if "csrf_token" not in session:
        session["csrf_token"] = secrets.token_hex(32)
    return jsonify(ok=True, csrf_token=session["csrf_token"])

# =============================================================================
# AUTH API ENDPOINTS
# =============================================================================
@app.route("/api/register", methods=["POST"])
def api_register():

    data = request.get_json() or {}

    name = data.get("name", "").strip()
    email = data.get("email", "").strip().lower()
    login_type = "email"

    if not name or not email:
        return jsonify(
            success=False,
            error="Name and Email are required."
        ), 400

    conn = sqlite3.connect(DATABASE)
    cur = conn.cursor()

    cur.execute(
        "SELECT id FROM users WHERE email=?",
        (email,)
    )

    if cur.fetchone():
        conn.close()
        return jsonify(
            success=False,
            error="Email already exists."
        ), 409

    cur.execute("""
    INSERT INTO users(
        name,
        email,
        login_type
    )
    VALUES(
        ?,
        ?,
        ?
    )
    """,
    (
        name,
        email,
        "email"
    )
)

    conn.commit()
    conn.close()

    return jsonify(
        success=True,
        message="Account created successfully."
    )

@app.route("/api/send-otp", methods=["POST"])
def send_otp():

    data = request.get_json()

    email = data.get("email","").strip()

    if not email:
        return jsonify({
            "ok": False,
            "message": "Email required"
        }),400

    otp = str(random.randint(100000,999999))
    print("OTP:", otp)
    print("EMAIL_ADDRESS:", os.getenv("EMAIL_ADDRESS"))
    print("EMAIL_PASSWORD:", os.getenv("EMAIL_PASSWORD"))

    otp_store[email] = {
        "otp": otp,
        "expires": datetime.utcnow() + timedelta(minutes=OTP_EXPIRE_MINUTES)
    }

    try:

        msg = EmailMessage()
        msg["Subject"] = "SYRA Login OTP"
        msg["From"] = os.getenv("EMAIL_ADDRESS")
        msg["To"] = email

        msg.set_content(f"""

Your SYRA Login OTP

OTP : {otp}

Valid for 5 minutes.

Do not share this OTP.

""")

        context = ssl.create_default_context()

        with smtplib.SMTP_SSL("smtp.gmail.com",465,context=context) as smtp:

            smtp.login(
                os.getenv("EMAIL_ADDRESS"),
                os.getenv("EMAIL_PASSWORD")
            )

            smtp.send_message(msg)
            print("Email sent successfully to:", email)

        return jsonify({
            "ok":True,
            "message":"OTP sent successfully"
        })

    except Exception as e:

        return jsonify({
            "ok":False,
            "message":str(e)
        }),500
    
@app.route("/api/verify-otp", methods=["POST"])
def verify_otp():

    data = request.get_json()

    email = data.get("email", "").strip().lower()
    otp = data.get("otp", "").strip()

    if not email or not otp:
        return jsonify({
            "ok": False,
            "message": "Email and OTP required"
        }), 400

    conn = sqlite3.connect(DATABASE)
    cur = conn.cursor()

    cur.execute(
        "SELECT name FROM users WHERE email=?",
        (email,)
    )

    row = cur.fetchone()

    if row is None:

        cur.execute(
            """
            INSERT INTO users(name,email,login_type)
            VALUES(?,?,?)
            """,
            (
                "User",
                email,
                "email"
            )
        )

        conn.commit()

        name = "User"

    else:

        name = row[0]

    conn.close()

    if email not in otp_store:
        return jsonify({
            "ok": False,
            "message": "OTP expired"
        }), 401

    saved = otp_store[email]

    if datetime.utcnow() > saved["expires"]:

        del otp_store[email]

        return jsonify({
            "ok": False,
            "message": "OTP expired"
        }), 401

    if otp != saved["otp"]:
        return jsonify({
            "ok": False,
            "message": "Invalid OTP"
        }), 401

    del otp_store[email]

    session["user"] = {
        "email": email,
        "name": name
    }

    return jsonify({
        "ok": True,
        "message": "Login successful",
        "user": {
            "email": email,
            "name": name
        }
    })

@app.route("/api/google-login", methods=["POST"])
def api_google_login():

    data = request.get_json() or {}

    id_token = data.get("idToken", "").strip()

    if not id_token:
        return jsonify({
            "ok": False,
            "message": "idToken is required"
        }), 400

    try:
        decoded_token = firebase_auth.verify_id_token(id_token)
    except Exception as e:
        return jsonify({
            "ok": False,
            "message": f"Invalid or expired token: {e}"
        }), 401

    uid = decoded_token.get("uid")
    email = (decoded_token.get("email") or "").strip().lower()
    name = decoded_token.get("name") or (email.split("@")[0] if email else "User")
    picture = decoded_token.get("picture")  # extracted per spec; no `picture`
                                             # column exists on users yet, so
                                             # it isn't persisted — add one if
                                             # you want it stored.

    if not email:
        return jsonify({
            "ok": False,
            "message": "Google account has no email"
        }), 401

    conn = sqlite3.connect(DATABASE)
    cur = conn.cursor()

    cur.execute(
        "SELECT id, name, google_id FROM users WHERE email=?",
        (email,)
    )

    row = cur.fetchone()

    if row is None:

        cur.execute(
            """
            INSERT INTO users(name, email, login_type, google_id, picture)
            VALUES(?, ?, ?, ?, ?)
            """,
            (
                name,
                email,
                "google",
                uid,
                picture
            )
        )

        conn.commit()

    else:

        existing_name = row[1]
        existing_google_id = row[2]

        if name != existing_name or uid != existing_google_id:

            cur.execute(
                """
                UPDATE users
                SET name = ?, google_id = ?
                WHERE email = ?
                """,
                (
                    name,
                    uid,
                    email
                )
            )

            conn.commit()

        # Keep the stored Google avatar in sync even when name/google_id
        # didn't change (e.g. the person updated their Google photo).
        if picture:
            cur.execute("UPDATE users SET picture=? WHERE email=?", (picture, email))
            conn.commit()

    conn.close()

    session["user"] = {
        "email": email,
        "name": name
    }
    session["csrf_token"] = session.get("csrf_token") or secrets.token_hex(32)

    return jsonify({
        "ok": True,
        "message": "Login successful",
        "user": {
            "name": name,
            "email": email
        }
    })


@app.route("/api/phone-login", methods=["POST"])
def api_phone_login():

    data = request.get_json() or {}

    id_token = data.get("idToken", "").strip()

    if not id_token:
        return jsonify({
            "ok": False,
            "message": "idToken is required"
        }), 400

    try:
        decoded_token = firebase_auth.verify_id_token(id_token)
    except Exception as e:
        return jsonify({
            "ok": False,
            "message": f"Invalid or expired token: {e}"
        }), 401

    uid = decoded_token.get("uid")
    phone = (decoded_token.get("phone_number") or "").strip()

    if not phone:
        return jsonify({
            "ok": False,
            "message": "Firebase token has no phone number"
        }), 401

    conn = sqlite3.connect(DATABASE)
    cur = conn.cursor()

    cur.execute(
        "SELECT id, phone FROM users WHERE phone=?",
        (phone,)
    )

    row = cur.fetchone()

    if row is None:

        cur.execute(
            """
            INSERT INTO users(name, phone, login_type, google_id)
            VALUES(?, ?, ?, ?)
            """,
            (
                "User",
                phone,
                "phone",
                uid
            )
        )

        conn.commit()

    else:

        existing_phone = row[1]

        if phone != existing_phone:

            cur.execute(
                """
                UPDATE users
                SET phone = ?
                WHERE id = ?
                """,
                (
                    phone,
                    row[0]
                )
            )

            conn.commit()

    conn.close()

    session["user"] = {
        "phone": phone,
        "name": "User"
    }
    session["csrf_token"] = session.get("csrf_token") or secrets.token_hex(32)

    return jsonify({
        "ok": True,
        "message": "Login successful",
        "user": {
            "phone": phone,
            "name": "User"
        }
    })


@app.route("/api/logout", methods=["POST"])
def api_logout():

    session.clear()

    return jsonify({
        "success": True
    })


@app.route("/api/check_session", methods=["GET"])
def api_check_session():

    user = session.get("user")

    if not user:
        return jsonify({
            "logged_in": False
        })

    return jsonify({
        "logged_in": True,
        "user": user
    })

# =============================================================================
# PROFILE (name, email, login provider, joined date, picture)
# =============================================================================
def _current_user_row(cur, user_key):
    """Looks up the users row matching the current session, trying email
    first (Google/Email-OTP login) then phone (phone login)."""
    row = cur.execute("SELECT * FROM users WHERE email=?", (user_key,)).fetchone()
    if row:
        return row
    return cur.execute("SELECT * FROM users WHERE phone=?", (user_key,)).fetchone()

@app.route("/api/profile", methods=["GET"])
def api_profile_get():
    try:
        user = session.get("user")
        if not user:
            return jsonify(ok=False, message="Not logged in"), 401

        user_key = resolve_session_id({})
        conn = get_db()
        try:
            row = _current_user_row(conn.cursor(), user_key)
        finally:
            conn.close()

        if not row:
            # Session exists but no DB row yet (e.g. an OTP-only session) —
            # fall back to whatever is cached in the Flask session.
            return jsonify(ok=True, profile={
                "name": user.get("name", "User"),
                "email": user.get("email", ""),
                "phone": user.get("phone", ""),
                "login_type": "email" if user.get("email") else ("phone" if user.get("phone") else ""),
                "picture": None,
                "joined_at": None,
            })

        return jsonify(ok=True, profile={
            "name": row["name"],
            "email": row["email"],
            "phone": row["phone"],
            "login_type": row["login_type"],
            "picture": row["picture"],
            "joined_at": row["created_at"],
        })
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500

@app.route("/api/profile", methods=["PATCH"])
def api_profile_update():
    try:
        user = session.get("user")
        if not user:
            return jsonify(ok=False, message="Not logged in"), 401

        req = request.get_json() or {}
        new_name = (req.get("name") or "").strip()[:120]
        if not new_name:
            return jsonify(ok=False, message="Name is required"), 400

        user_key = resolve_session_id({})
        conn = get_db()
        try:
            cur = conn.cursor()
            if "@" in user_key:
                cur.execute("UPDATE users SET name=? WHERE email=?", (new_name, user_key))
            else:
                cur.execute("UPDATE users SET name=? WHERE phone=?", (new_name, user_key))
            conn.commit()
        finally:
            conn.close()

        session["user"]["name"] = new_name
        return jsonify(ok=True, profile={"name": new_name})
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500

@app.route("/api/profile/picture", methods=["POST"])
def api_profile_picture():
    try:
        user = session.get("user")
        if not user:
            return jsonify(ok=False, message="Not logged in"), 401

        if "file" not in request.files:
            return jsonify(ok=False, message="No file part in request"), 400
        file = request.files["file"]
        if not file or file.filename == "":
            return jsonify(ok=False, message="No file selected"), 400

        ext = file.filename.rsplit(".", 1)[-1].lower() if "." in file.filename else ""
        if ext not in IMAGE_EXTENSIONS:
            return jsonify(ok=False, message="Unsupported image type"), 400

        user_key = resolve_session_id({})
        safe_name = secure_filename(user_key) or "user"
        filename = f"{safe_name}_{secrets.token_hex(6)}.{ext}"
        save_path = os.path.join(AVATAR_FOLDER, filename)
        file.save(save_path)
        picture_url = f"/avatars/{filename}"

        conn = get_db()
        try:
            cur = conn.cursor()
            if "@" in user_key:
                cur.execute("UPDATE users SET picture=? WHERE email=?", (picture_url, user_key))
            else:
                cur.execute("UPDATE users SET picture=? WHERE phone=?", (picture_url, user_key))
            conn.commit()
        finally:
            conn.close()

        return jsonify(ok=True, picture=picture_url)
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500

@app.route("/avatars/<path:filename>")
def serve_avatar(filename):
    return send_from_directory(AVATAR_FOLDER, filename)

# =============================================================================
# SETTINGS (dark mode, language, theme)
# =============================================================================
@app.route("/api/settings", methods=["GET"])
def api_settings_get():
    try:
        user_key = resolve_session_id({"session_id": request.args.get("session_id")})
        conn = get_db()
        try:
            row = conn.execute(
                "SELECT * FROM user_settings WHERE user_key=?", (user_key,)
            ).fetchone()
        finally:
            conn.close()

        if not row:
            return jsonify(ok=True, settings={"dark_mode": True, "language": "en", "theme": "dark"})

        return jsonify(ok=True, settings={
            "dark_mode": bool(row["dark_mode"]),
            "language": row["language"],
            "theme": row["theme"],
        })
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500

@app.route("/api/settings", methods=["PATCH"])
def api_settings_update():
    try:
        req = request.get_json() or {}
        user_key = resolve_session_id(req)

        conn = get_db()
        try:
            cur = conn.cursor()
            existing = cur.execute(
                "SELECT * FROM user_settings WHERE user_key=?", (user_key,)
            ).fetchone()

            dark_mode = req.get("dark_mode", bool(existing["dark_mode"]) if existing else True)
            language = (req.get("language") or (existing["language"] if existing else "en"))[:10]
            theme = (req.get("theme") or (existing["theme"] if existing else "dark"))[:20]
            now = _now_iso()

            if existing:
                cur.execute(
                    "UPDATE user_settings SET dark_mode=?, language=?, theme=?, updated_at=? WHERE user_key=?",
                    (1 if dark_mode else 0, language, theme, now, user_key)
                )
            else:
                cur.execute(
                    "INSERT INTO user_settings(user_key, dark_mode, language, theme, updated_at) VALUES (?,?,?,?,?)",
                    (user_key, 1 if dark_mode else 0, language, theme, now)
                )
            conn.commit()
        finally:
            conn.close()

        return jsonify(ok=True, settings={"dark_mode": bool(dark_mode), "language": language, "theme": theme})
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500

# =============================================================================
# CHAT HISTORY: CLEAR ALL / EXPORT
# =============================================================================
@app.route("/api/chats/clear", methods=["POST"])
def api_chats_clear():
    try:
        req = request.get_json() or {}
        user_key = resolve_session_id(req)
        conn = get_db()
        try:
            cur = conn.cursor()
            cur.execute("DELETE FROM messages WHERE user_key=?", (user_key,))
            cur.execute("DELETE FROM conversations WHERE user_key=?", (user_key,))
            conn.commit()
        finally:
            conn.close()
        _histories.clear()
        return jsonify(ok=True, message="Chat history cleared")
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500

@app.route("/api/chats/export", methods=["GET"])
def api_chats_export():
    try:
        user_key = resolve_session_id({"session_id": request.args.get("session_id")})
        conn = get_db()
        try:
            convs = conn.execute(
                "SELECT * FROM conversations WHERE user_key=? ORDER BY updated_at DESC", (user_key,)
            ).fetchall()
            export = []
            for c in convs:
                msgs = conn.execute(
                    "SELECT role, content, created_at FROM messages WHERE user_key=? AND conv_id=? ORDER BY id ASC",
                    (user_key, c["id"])
                ).fetchall()
                export.append({
                    "id": c["id"],
                    "title": c["title"],
                    "pinned": bool(c["pinned"]),
                    "created_at": c["created_at"],
                    "updated_at": c["updated_at"],
                    "messages": [{"role": m["role"], "content": m["content"], "created_at": m["created_at"]} for m in msgs],
                })
        finally:
            conn.close()

        payload = json.dumps({"exported_at": _now_iso(), "conversations": export}, indent=2)
        return Response(
            payload,
            mimetype="application/json",
            headers={"Content-Disposition": "attachment; filename=syra_chat_export.json"}
        )
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500

# =============================================================================
# ACCOUNT DELETION
# =============================================================================
@app.route("/api/account/delete", methods=["POST"])
def api_account_delete():
    try:
        user = session.get("user")
        if not user:
            return jsonify(ok=False, message="Not logged in"), 401

        user_key = resolve_session_id({})
        conn = get_db()
        try:
            cur = conn.cursor()
            cur.execute("DELETE FROM messages WHERE user_key=?", (user_key,))
            cur.execute("DELETE FROM conversations WHERE user_key=?", (user_key,))
            cur.execute("DELETE FROM user_settings WHERE user_key=?", (user_key,))
            if "@" in user_key:
                cur.execute("DELETE FROM users WHERE email=?", (user_key,))
            else:
                cur.execute("DELETE FROM users WHERE phone=?", (user_key,))
            conn.commit()
        finally:
            conn.close()

        _histories.clear()
        session.clear()
        return jsonify(ok=True, message="Account deleted")
    except Exception as e:
        traceback.print_exc()
        return jsonify(ok=False, message=str(e)), 500

# =============================================================================
# STATIC SITES & HOME
# =============================================================================
@app.route("/")
def serve_index():
    return send_from_directory(PUBLIC_FOLDER,"index.html")

@app.route("/<path:path>")
def serve_static(path):
    return send_from_directory(PUBLIC_FOLDER,path)

print("SYRA Backend Started")

print(app.url_map)

init_database()

if __name__ == "__main__":
    # Production readiness: debug mode is now opt-in via FLASK_DEBUG (or
    # FLASK_ENV=development), instead of being hardcoded on. Development
    # still works the same as before when you explicitly enable it.
    _debug = os.environ.get("FLASK_DEBUG", "").lower() in ("1", "true", "yes") or \
             os.environ.get("FLASK_ENV", "").lower() == "development"
    _port = int(os.environ.get("PORT", 5000))
    app.run(host="0.0.0.0", port=_port, debug=_debug)
